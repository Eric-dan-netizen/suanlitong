"""
算力通 · Pitch Deck 生成器
种子轮融资路演 PPTX 自动生成
"""


from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

# ── Global Design Tokens ──────────────────────────────
W = Inches(13.333)  # 16:9 widescreen
H = Inches(7.5)

BG_DARK   = RGBColor(0x0A, 0x0E, 0x17)   # deep navy
BG_CARD   = RGBColor(0x12, 0x18, 0x24)   # card / panel bg
ACCENT    = RGBColor(0x00, 0xD4, 0xFF)   # cyan highlight
ACCENT2   = RGBColor(0x00, 0xFF, 0x88)   # green data
ACCENT3   = RGBColor(0xFF, 0x6B, 0x35)   # orange warning
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GRAY      = RGBColor(0x88, 0x93, 0xA0)
GRAY_DIM  = RGBColor(0x55, 0x5E, 0x6B)
BORDER    = RGBColor(0x1E, 0x28, 0x36)

FONT_TITLE = "PingFang SC"
FONT_BODY  = "PingFang SC"

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

# Use blank layout
blank_layout = prs.slide_layouts[6]  # blank


# ── Helper Functions ──────────────────────────────────

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color=None, border_color=None, border_width=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        left, top, width, height
    )
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.fill.solid()
        if border_width:
            shape.line.width = border_width
    return shape

def add_line(slide, left, top, width, height, color, line_width=Pt(1)):
    """Add a simple horizontal line using a thin rectangle."""
    return add_rect(slide, left, top, width, height, fill_color=color)

def add_textbox(slide, left, top, width, height, text="", font_size=Pt(14),
                font_color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                font_name=FONT_BODY, anchor=MSO_ANCHOR.TOP, line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    # Set anchor
    txBox.text_frame.paragraphs[0].alignment = alignment

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.line_spacing = Pt(int(font_size.pt * line_spacing))

    # Set East Asian font
    for run in p.runs:
        rPr = run._r.get_or_add_rPr()
        rPr.set(qn('a:altLang'), 'zh-CN')

    return txBox

def add_multiline_textbox(slide, left, top, width, height, lines, default_size=Pt(14),
                          default_color=WHITE, default_bold=False, alignment=PP_ALIGN.LEFT,
                          font_name=FONT_BODY, anchor=MSO_ANCHOR.TOP, line_spacing=1.3):
    """lines: list of (text, font_size, font_color, bold) or just str"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if isinstance(line, str):
            p.text = line
            p.font.size = default_size
            p.font.color.rgb = default_color
            p.font.bold = default_bold
        else:
            text, size, color, bold = line[0], line[1] if len(line) > 1 else default_size, \
                                      line[2] if len(line) > 2 else default_color, \
                                      line[3] if len(line) > 3 else default_bold
            p.text = text
            p.font.size = size
            p.font.color.rgb = color
            p.font.bold = bold

        p.font.name = font_name
        p.alignment = alignment
        p.line_spacing = Pt(int((size if isinstance(line, tuple) and len(line) > 1 else default_size).pt * line_spacing))

    return txBox

def add_card(slide, left, top, width, height, title="", body_lines=None,
             title_color=ACCENT, body_color=WHITE, title_size=Pt(16), body_size=Pt(12),
             fill=BG_CARD, border=BORDER):
    """Add a card-style panel with title and body text."""
    card = add_rect(slide, left, top, width, height, fill_color=fill, border_color=border, border_width=Pt(0.5))

    y_offset = top + Inches(0.15)
    if title:
        add_textbox(slide, left + Inches(0.2), y_offset, width - Inches(0.4), Inches(0.3),
                    text=title, font_size=title_size, font_color=title_color, bold=True)
        y_offset += Inches(0.35)

    if body_lines:
        add_multiline_textbox(slide, left + Inches(0.2), y_offset, width - Inches(0.4),
                              height - (y_offset - top) - Inches(0.15),
                              lines=body_lines, default_size=body_size, default_color=body_color)

    return card

def add_big_number(slide, left, top, number_text, label_text, num_color=ACCENT2, label_color=WHITE):
    """Add a large number with label below it."""
    add_textbox(slide, left, top, Inches(3), Inches(0.8),
                text=number_text, font_size=Pt(42), font_color=num_color, bold=True)
    add_textbox(slide, left, top + Inches(0.75), Inches(3), Inches(0.4),
                text=label_text, font_size=Pt(14), font_color=label_color, bold=False)

def add_slide_number(slide, num):
    add_textbox(slide, W - Inches(0.8), H - Inches(0.4), Inches(0.6), Inches(0.3),
                text=str(num), font_size=Pt(10), font_color=GRAY_DIM, alignment=PP_ALIGN.RIGHT)

def add_section_title(slide, title, subtitle=None):
    """Top-of-slide section title with accent underline."""
    add_textbox(slide, Inches(0.6), Inches(0.35), Inches(12), Inches(0.5),
                text=title, font_size=Pt(28), font_color=WHITE, bold=True)
    add_line(slide, Inches(0.6), Inches(0.9), Inches(1.2), Pt(3), ACCENT)
    if subtitle:
        add_textbox(slide, Inches(0.6), Inches(1.05), Inches(12), Inches(0.35),
                    text=subtitle, font_size=Pt(13), font_color=GRAY)

def add_footer_note(slide, text):
    add_textbox(slide, Inches(0.6), H - Inches(0.4), Inches(9), Inches(0.3),
                text=text, font_size=Pt(9), font_color=GRAY_DIM)


# ══════════════════════════════════════════════════════
# SLIDE 1 · COVER
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG_DARK)

# Top decorative line
add_line(slide, Inches(0), Inches(0), W, Pt(4), ACCENT)

# Center content
add_textbox(slide, Inches(1.5), Inches(1.8), Inches(10.3), Inches(1.2),
            text="算 力 通", font_size=Pt(60), font_color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1.5), Inches(2.9), Inches(10.3), Inches(0.6),
            text="SuanLiTong", font_size=Pt(22), font_color=ACCENT,
            alignment=PP_ALIGN.CENTER)

# Accent divider
add_line(slide, Inches(5.2), Inches(3.55), Inches(3), Pt(2), ACCENT)

add_textbox(slide, Inches(1.5), Inches(3.8), Inches(10.3), Inches(0.6),
            text="让算力像自来水一样，扭开即用", font_size=Pt(18), font_color=GRAY,
            alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1.5), Inches(4.6), Inches(10.3), Inches(0.5),
            text="GPU 算力多云聚合平台 · 种子轮融资", font_size=Pt(16), font_color=GRAY,
            alignment=PP_ALIGN.CENTER)

# Bottom info
add_textbox(slide, Inches(1.5), Inches(5.8), Inches(10.3), Inches(0.4),
            text="创始人：Eric  |  算力通  |  2026.07", font_size=Pt(12), font_color=GRAY_DIM,
            alignment=PP_ALIGN.CENTER)

# Bottom decorative line
add_line(slide, Inches(0), H - Pt(4), W, Pt(4), ACCENT)


# ══════════════════════════════════════════════════════
# SLIDE 2 · 供需缺口
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG_DARK)
add_slide_number(slide, 2)

add_section_title(slide, "一个正在爆炸的供需缺口",
                  "2026 Q1 中国 AI 算力市场 · 数据来源：工信部、格隆汇、阿里云官方公告")

# Big numbers side by side
add_big_number(slide, Inches(1.2), Inches(1.8), "+417%", "需求增速（同比）")
add_big_number(slide, Inches(4.8), Inches(1.8), "VS", "", num_color=GRAY)
add_big_number(slide, Inches(6.2), Inches(1.8), "+128%", "供给增速（同比）")

# Gap visualization bar
add_rect(slide, Inches(1.2), Inches(3.1), Inches(10.9), Inches(0.5), fill_color=BG_CARD, border_color=BORDER)
# Demand bar (wider)
add_rect(slide, Inches(1.2), Inches(3.15), Inches(8.8), Inches(0.4), fill_color=ACCENT3)
# Supply bar (narrower, overlaid)
add_rect(slide, Inches(1.2), Inches(3.7), Inches(2.9), Inches(0.4), fill_color=ACCENT)

add_textbox(slide, Inches(1.2), Inches(3.85), Inches(4), Inches(0.3),
            text="需求", font_size=Pt(11), font_color=ACCENT3)
add_textbox(slide, Inches(1.2), Inches(4.05), Inches(4), Inches(0.3),
            text="供给", font_size=Pt(11), font_color=ACCENT)

# Key facts in cards
facts = [
    ("H100 排期至 2028 年", "高端 GPU 供不应求，全球排队"),
    ("GPU 租赁率 > 90%", "上架即秒光，行业常态"),
    ("阿里云 GPU 涨价 +20%~34%", "2026 年 3 月官方调价"),
]
for i, (title, desc) in enumerate(facts):
    x = Inches(1.2 + i * 3.8)
    add_card(slide, x, Inches(4.5), Inches(3.5), Inches(1.2),
             title=title, body_lines=[desc],
             title_size=Pt(15), body_size=Pt(11))

# Bottom punchline
add_textbox(slide, Inches(1.2), Inches(6.0), Inches(10.9), Inches(0.5),
            text="每 5 份算力需求，市场只接得住 2 份。缺口仍在扩大。",
            font_size=Pt(16), font_color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════
# SLIDE 3 · 用户痛点
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG_DARK)
add_slide_number(slide, 3)

add_section_title(slide, "940 万 AI 开发者的真实困境",
                  "数据来源：知乎 × 魔搭《2025 AI 开发者生态白皮书》")

# Top stats
add_big_number(slide, Inches(1.2), Inches(1.6), "940 万", "中国 AI 开发者总量", num_color=ACCENT2)
add_big_number(slide, Inches(4.8), Inches(1.6), "129 万", "独立开发者", num_color=ACCENT2)
add_big_number(slide, Inches(8.4), Inches(1.6), "195 万", "50人以下小团队", num_color=ACCENT2)

# Pain points as numbered cards
pains = [
    ("01", "买不起", "H100 一张 20 万\n一台 8 卡 200 万起"),
    ("02", "配不好", "阿里云后台配环境\n新手需要 2 天"),
    ("03", "比不到", "同一张卡不同云厂\n差价 20% ~ 40%"),
    ("04", "锁太久", "包月起售\n只用 50 小时也付整月"),
    ("05", "重装烦", "每次新建实例\n重装 CUDA / PyTorch"),
]
for i, (num, title, desc) in enumerate(pains):
    x = Inches(0.6 + i * 2.5)
    # Number circle
    add_rect(slide, x + Inches(0.7), Inches(2.9), Inches(0.6), Inches(0.6),
             fill_color=ACCENT, border_color=None)
    add_textbox(slide, x + Inches(0.7), Inches(2.92), Inches(0.6), Inches(0.55),
                text=num, font_size=Pt(20), font_color=BG_DARK, bold=True, alignment=PP_ALIGN.CENTER)
    # Title
    add_textbox(slide, x + Inches(0.1), Inches(3.65), Inches(2.1), Inches(0.35),
                text=title, font_size=Pt(16), font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Desc
    add_textbox(slide, x + Inches(0.1), Inches(4.05), Inches(2.1), Inches(0.9),
                text=desc, font_size=Pt(12), font_color=GRAY, alignment=PP_ALIGN.CENTER)

# Bottom
add_textbox(slide, Inches(1.2), Inches(5.6), Inches(10.9), Inches(0.5),
            text="他们要的不是更便宜的 GPU，是“打开浏览器就能写代码”。",
            font_size=Pt(16), font_color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════
# SLIDE 4 · 解决方案
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG_DARK)
add_slide_number(slide, 4)

add_section_title(slide, "算力通：多云比价，一键即用",
                  "第一阶段：0 张自有 GPU，100% 多云代理商模式")

# 6 feature cards in 3x2 grid
features = [
    ("🔍 多云比价", "阿里云 · 华为云 · 腾讯云\n实时比价，自动推荐最优"),
    ("⚡ 一键部署", "选 PyTorch → 3 分钟就绪\n不用配 VPC / CUDA / NAS"),
    ("⏱ 按秒计费", "用 37 分钟付 37 分钟\n不花一分冤枉钱"),
    ("📦 50+ 镜像", "预置主流 AI 框架\n开箱即用，零配置"),
    ("💸 Spot 竞价", "凌晨闲置 GPU 最低 3 折\n平台自动匹配低价资源"),
    ("📋 算力券代申领", "政府补贴帮你搞定\n贵州补贴 30%，内蒙年 5000 万"),
]

card_w = Inches(3.8)
card_h = Inches(1.45)
gap_x = Inches(0.3)
gap_y = Inches(0.25)
start_x = Inches(0.6)
start_y = Inches(1.6)

for i, (title, desc) in enumerate(features):
    col = i % 3
    row = i // 3
    x = start_x + col * (card_w + gap_x)
    y = start_y + row * (card_h + gap_y)

    card = add_rect(slide, x, y, card_w, card_h, fill_color=BG_CARD, border_color=BORDER, border_width=Pt(0.5))
    add_textbox(slide, x + Inches(0.2), y + Inches(0.15), card_w - Inches(0.4), Inches(0.3),
                text=title, font_size=Pt(15), font_color=ACCENT, bold=True)
    add_textbox(slide, x + Inches(0.2), y + Inches(0.55), card_w - Inches(0.4), Inches(0.8),
                text=desc, font_size=Pt(11), font_color=GRAY)

# Bottom tagline
add_textbox(slide, Inches(1.2), Inches(4.85), Inches(10.9), Inches(0.5),
            text="你选框架，我们搞定 GPU。", font_size=Pt(18), font_color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════
# SLIDE 5 · WHY NOW
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG_DARK)
add_slide_number(slide, 5)

add_section_title(slide, "为什么是现在？",
                  "四个不可复制的时机窗口 · 窗口预计到 2028 年国产 GPU 量产开始收窄")

reasons = [
    ("供需缺口历史最大", "需求 +417% vs 供给 +128%\n2028 年国产 GPU 量产前不会收窄",
     ACCENT3),
    ("政策历史级加码", "“十五五”首次将算力网与\n水网、电网并列为国家基础设施\n地方算力券密集推出",
     ACCENT2),
    ("开发者“超级个体”化", "AI 编程工具让个人 = 一个团队\n129 万独立开发者的算力需求\n大云厂商看不上（收入占比 <5%）",
     ACCENT),
    ("同类项目已验证融资路径", "密瓜智能 种子轮 500万+\n→ 天使轮数千万\n赛道融资热度持续",
     RGBColor(0xFF, 0xD7, 0x00)),
]

for i, (title, desc, color) in enumerate(reasons):
    x = Inches(0.6 + i * 3.2)
    # Accent top line
    add_rect(slide, x, Inches(1.6), Inches(2.9), Pt(4), fill_color=color)
    add_card(slide, x, Inches(1.75), Inches(2.9), Inches(2.5),
             title=title, body_lines=[desc],
             title_color=color, title_size=Pt(16), body_size=Pt(12))


# ══════════════════════════════════════════════════════
# SLIDE 6 · 市场空间
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG_DARK)
add_slide_number(slide, 6)

add_section_title(slide, "市场空间——自下而上算的账",
                  "不拍脑袋，从真实用户数往上乘")

# TAM/SAM/SOM flow
add_textbox(slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(0.4),
            text="目标用户池  324 万人  （195 万小团队 + 129 万独立开发者）",
            font_size=Pt(15), font_color=WHITE)

add_textbox(slide, Inches(0.8), Inches(2.1), Inches(11.5), Inches(0.4),
            text="↓  需 GPU 的 AI 开发者（60%）  →  194 万人",
            font_size=Pt(15), font_color=GRAY)

add_textbox(slide, Inches(0.8), Inches(2.6), Inches(11.5), Inches(0.4),
            text="↓  × 人均年 GPU 消费 5,000 元（参考 AutoDL 用户行为数据）",
            font_size=Pt(15), font_color=GRAY)

add_textbox(slide, Inches(0.8), Inches(3.1), Inches(11.5), Inches(0.4),
            text="↓",
            font_size=Pt(15), font_color=GRAY)

# Three columns: TAM SAM SOM
cols = [
    ("TAM", "≈ 97 亿元/年", "总可触达市场", ACCENT2),
    ("SAM", "≈ 50 亿元/年", "可服务市场（聚焦创业公司\n+独立开发者）", ACCENT),
    ("SOM", "2,500万 ~ 1亿/年", "3年后市场份额 0.5%~2%", RGBColor(0xFF, 0xD7, 0x00)),
]
for i, (label, amount, desc, color) in enumerate(cols):
    x = Inches(0.8 + i * 4.1)
    add_rect(slide, x, Inches(3.6), Inches(3.8), Inches(2.1), fill_color=BG_CARD, border_color=color, border_width=Pt(1))
    add_textbox(slide, x + Inches(0.25), Inches(3.75), Inches(3.3), Inches(0.35),
                text=label, font_size=Pt(14), font_color=color, bold=True)
    add_textbox(slide, x + Inches(0.25), Inches(4.1), Inches(3.3), Inches(0.5),
                text=amount, font_size=Pt(26), font_color=WHITE, bold=True)
    add_textbox(slide, x + Inches(0.25), Inches(4.7), Inches(3.3), Inches(0.7),
                text=desc, font_size=Pt(12), font_color=GRAY)

# Global context
add_textbox(slide, Inches(0.8), Inches(6.1), Inches(11.5), Inches(0.4),
            text="全球 GPU 算力租赁市场：74 亿美元（2025）→ 282 亿美元（2032）· CAGR 21.4%",
            font_size=Pt(13), font_color=GRAY)
add_footer_note(slide, "数据来源：知乎×魔搭白皮书、AutoDL 公开定价、格隆汇、辰宇信息")


# ══════════════════════════════════════════════════════
# SLIDE 7 · 产品三阶段
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG_DARK)
add_slide_number(slide, 7)

add_section_title(slide, "产品——三阶段演进路线",
                  "核心逻辑：用代理验证 PMF → 用 PMF 融资 → 用资本建资产")

phases = [
    ("Phase 1 · 纯代理", "0 ~ 12 月", "毛利率 10~25%",
     "多云聚合，零资产启动。\n阿里云 + 华为云 + 腾讯云\nGPU 实例统一调度。",
     ACCENT),
    ("Phase 2 · 混合云", "12 ~ 24 月", "毛利率 25~45%",
     "自有 GPU + 云 Spot + 闲置撮合\n三轨并行。首批 8~16 张 H100\n部署西部枢纽（低电价 + 算力券）。",
     ACCENT2),
    ("Phase 3 · 资产运营", "24 ~ 36 月+", "毛利率 40~55%",
     "自有 GPU 规模化。\n算力期货/远期合约。\n调度 SaaS 输出给 IDC 运营商。",
     RGBColor(0xFF, 0xD7, 0x00)),
]

for i, (title, period, margin, desc, color) in enumerate(phases):
    x = Inches(0.6 + i * 4.2)
    # Phase header
    add_rect(slide, x, Inches(1.6), Inches(3.9), Inches(0.5), fill_color=color)
    add_textbox(slide, x + Inches(0.15), Inches(1.62), Inches(3.6), Inches(0.45),
                text=f"{title}    {period}    {margin}", font_size=Pt(13), font_color=BG_DARK, bold=True)

    # Body
    add_card(slide, x, Inches(2.2), Inches(3.9), Inches(2.5),
             title="", body_lines=[desc], body_size=Pt(13))

    # Arrow between phases (except last)
    if i < 2:
        add_textbox(slide, x + Inches(3.95), Inches(3.0), Inches(0.4), Inches(0.4),
                    text="→", font_size=Pt(24), font_color=ACCENT, bold=True)


# ══════════════════════════════════════════════════════
# SLIDE 8 · 盈利模型
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG_DARK)
add_slide_number(slide, 8)

add_section_title(slide, "盈利模型——不是讲故事，是算账",
                  "假设来源：AutoDL 用户行为数据、云厂商代理政策、IDC 市场报价")

# Left: unit economics
add_textbox(slide, Inches(0.6), Inches(1.5), Inches(5.5), Inches(0.35),
            text="代理模式单位经济学（第一阶段）", font_size=Pt(17), font_color=ACCENT, bold=True)

ue_card = add_rect(slide, Inches(0.6), Inches(1.95), Inches(5.8), Inches(3.6),
                   fill_color=BG_CARD, border_color=BORDER, border_width=Pt(0.5))

ue_data = [
    ("CAC（用户获取成本）",     "150 元"),
    ("ARPU（月均消费）",        "600 元/月"),
    ("毛利率",                  "20%"),
    ("用户平均生命周期",         "20 个月"),
    ("", ""),
    ("LTV = 600 × 20% × 20",  "= 2,400 元"),
]
for i, (label, value) in enumerate(ue_data):
    y = Inches(2.1 + i * 0.42)
    add_textbox(slide, Inches(0.85), y, Inches(3.0), Inches(0.35),
                text=label, font_size=Pt(13), font_color=GRAY)
    add_textbox(slide, Inches(4.0), y, Inches(2.2), Inches(0.35),
                text=value, font_size=Pt(13), font_color=WHITE, bold=True, alignment=PP_ALIGN.RIGHT)

# Highlight box: LTV:CAC
highlight = add_rect(slide, Inches(0.85), Inches(4.95), Inches(5.3), Inches(0.5),
                     fill_color=RGBColor(0x0A, 0x30, 0x20), border_color=ACCENT2, border_width=Pt(1))
add_textbox(slide, Inches(1.0), Inches(4.98), Inches(5.0), Inches(0.44),
            text="LTV : CAC = 16 : 1    （SaaS 健康基准 ≥ 3:1）   回本周期：1.25 个月",
            font_size=Pt(14), font_color=ACCENT2, bold=True)

# Right: self-built GPU
add_textbox(slide, Inches(7.2), Inches(1.5), Inches(5.5), Inches(0.35),
            text="自有 GPU 经济账（单台 8 卡 H100）", font_size=Pt(17), font_color=ACCENT2, bold=True)

gpu_card = add_rect(slide, Inches(7.2), Inches(1.95), Inches(5.8), Inches(3.6),
                    fill_color=BG_CARD, border_color=BORDER, border_width=Pt(0.5))

gpu_data = [
    ("投入（CAPEX）",            "164 万元"),
    ("月运营成本（OPEX）",        "6.6 万元"),
    ("月收入（70% 利用率 @35元/时）", "14.1 万元"),
    ("月毛利",                   "7.5 万元"),
    ("毛利率",                   "53%"),
    ("投资回收期",                "22 个月"),
]
for i, (label, value) in enumerate(gpu_data):
    y = Inches(2.1 + i * 0.47)
    add_textbox(slide, Inches(7.45), y, Inches(3.2), Inches(0.37),
                text=label, font_size=Pt(13), font_color=GRAY)
    add_textbox(slide, Inches(10.8), y, Inches(2.0), Inches(0.37),
                text=value, font_size=Pt(13), font_color=WHITE, bold=True, alignment=PP_ALIGN.RIGHT)

# Note
add_textbox(slide, Inches(7.45), Inches(5.0), Inches(5.3), Inches(0.4),
            text="西部枢纽部署（低电价 + 算力券 30%）可缩至 18 个月回本",
            font_size=Pt(11), font_color=GRAY)


# ══════════════════════════════════════════════════════
# SLIDE 9 · 竞品
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG_DARK)
add_slide_number(slide, 9)

add_section_title(slide, "竞争格局——结构性空白",
                  "大厂“看不上”个人开发者 · 小厂“够不到”企业级 GPU")

# Comparison table
table_data = [
    ["", "大云厂商\n阿里云 / 华为云", "算力通", "第三方平台\nAutoDL / 恒源云"],
    ["GPU 类型", "企业级为主", "企业级为主", "消费级为主\n(4090 / 3090)"],
    ["计费粒度", "按小时 / 包月", "按秒", "按小时"],
    ["跨云比价", "✗", "✓", "✗"],
    ["开箱即用", "✗（需配 VPC/NAS/CUDA）", "✓（3 分钟就绪）", "✓"],
    ["算力券代申领", "✗", "✓", "✗"],
    ["合规增值服务", "✗", "✓（法律+金融背景）", "✗"],
    ["目标用户", "政企 / 中大互联网", "创业公司\n独立开发者", "学生 / 个人开发者"],
]

col_w = [Inches(2.0), Inches(3.2), Inches(3.2), Inches(3.2)]
row_h = Inches(0.55)
table_left = Inches(0.7)
table_top = Inches(1.5)

for r, row in enumerate(table_data):
    x = table_left
    for c, cell_text in enumerate(row):
        is_header = (r == 0)
        is_highlight = (c == 2 and r > 0)

        bg = ACCENT if (is_header and c == 2) else (BG_CARD if is_header else None)
        if is_highlight:
            bg = RGBColor(0x0A, 0x25, 0x30)

        y = table_top + r * row_h
        rect = add_rect(slide, x, y, col_w[c], row_h,
                        fill_color=bg, border_color=BORDER, border_width=Pt(0.5))

        font_c = WHITE if (is_header and c == 2) else (ACCENT if is_highlight else WHITE)
        add_textbox(slide, x + Inches(0.15), y + Inches(0.05), col_w[c] - Inches(0.3), row_h - Inches(0.1),
                    text=cell_text, font_size=Pt(12 if r > 0 else 13), font_color=font_c,
                    bold=(r == 0), alignment=PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT)

        x += col_w[c]

# Bottom insight
add_textbox(slide, Inches(0.7), Inches(6.1), Inches(11.5), Inches(0.5),
            text="结构性空白：企业级 GPU + 开发者级体验 + 合规壁垒 → 大厂和小厂都复制不了的三重叠加",
            font_size=Pt(15), font_color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════
# SLIDE 10 · GTM
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG_DARK)
add_slide_number(slide, 10)

add_section_title(slide, "怎么获客——第一批 500 个付费用户",
                  "三条增长飞轮：开源引流 + 内容营销 + 政策杠杆")

channels = [
    ("开源引流", ACCENT,
     ["GitHub 开源 GPU 比价 CLI 工具",
      "Python SDK（pip install）",
      "从贡献者中筛选早期员工",
      "Star 数 → 自然流量 → 注册"]),
    ("内容营销", ACCENT2,
     ["“AI 创业算力成本优化”系列",
      "知乎 / CSDN / 掘金 / 即刻",
      "真实数据 + 方法论",
      "不投硬广，靠搜索长尾流量"]),
    ("算力券杠杆", RGBColor(0xFF, 0xD7, 0x00),
     ["免费帮开发者代申领算力券",
      "贵州 30% · 呼市 5000 万/年",
      "领到券 → 自然留在平台消费",
      "政府补贴变获客成本"]),
]

for i, (title, color, bullets) in enumerate(channels):
    x = Inches(0.6 + i * 4.2)
    add_rect(slide, x, Inches(1.6), Inches(3.9), Inches(0.45), fill_color=color)
    add_textbox(slide, x + Inches(0.15), Inches(1.62), Inches(3.6), Inches(0.4),
                text=title, font_size=Pt(15), font_color=BG_DARK, bold=True)

    add_card(slide, x, Inches(2.15), Inches(3.9), Inches(2.0),
             title="", body_lines=[(b, Pt(12), GRAY, False) for b in bullets],
             body_size=Pt(12))

# Funnel summary at bottom
add_rect(slide, Inches(0.6), Inches(4.5), Inches(12.1), Inches(0.65), fill_color=BG_CARD, border_color=BORDER)
funnel_items = [
    ("开源工具 + 内容 + 算力券", "→", "注册用户", "→", "免费 20 GPU 小时", "→", "付费转化", "→", "团队邀请"),
]
add_textbox(slide, Inches(0.8), Inches(4.55), Inches(11.7), Inches(0.55),
            text="开源工具 + 内容 + 算力券  →  注册用户  →  免费 20 GPU 小时体验  →  付费转化  →  邀请团队成员",
            font_size=Pt(14), font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Milestones
milestones_text = "M3: 1,000 注册 / 30 付费    →    M6: 3,000 注册 / 120 付费    →    M12: 10,000 注册 / 500 付费 / MRR 30 万"
add_textbox(slide, Inches(0.6), Inches(5.4), Inches(12.1), Inches(0.5),
            text=milestones_text, font_size=Pt(14), font_color=ACCENT2, bold=True, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════
# SLIDE 11 · 三年路线图
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG_DARK)
add_slide_number(slide, 11)

add_section_title(slide, "三年路线图与 A 轮触发条件",
                  "不靠“觉得差不多了”，靠具体的数据指标触发下一轮融资")

# Timeline - 4 milestones
milestones = [
    ("2026 Q3", "种子轮", "200 万到位\n3 人团队组建\nMVP 启动开发", ACCENT),
    ("2027 Q1", "MVP 上线", "阿里云 + 华为云对接\n1,000 注册 / 30 付费\n10 个预置镜像", ACCENT),
    ("2027 Q3", "MRR 30 万", "500 付费用户\n盈亏平衡\n启动 A 轮融资", ACCENT2),
    ("2028 Q2", "混合云上线", "首批自有 GPU 部署\nARR 960 万+\n多区域扩张", RGBColor(0xFF, 0xD7, 0x00)),
]

for i, (date, title, desc, color) in enumerate(milestones):
    x = Inches(0.6 + i * 3.2)
    # Timeline node
    add_rect(slide, x + Inches(1.2), Inches(1.6), Inches(0.5), Inches(0.5), fill_color=color)
    add_textbox(slide, x + Inches(1.2), Inches(1.62), Inches(0.5), Inches(0.46),
                text=str(i+1), font_size=Pt(18), font_color=BG_DARK, bold=True, alignment=PP_ALIGN.CENTER)
    # Connector line
    if i < 3:
        add_line(slide, x + Inches(1.7), Inches(1.83), Inches(2.5), Pt(3), color)

    add_textbox(slide, x + Inches(0.1), Inches(2.25), Inches(2.8), Inches(0.3),
                text=date, font_size=Pt(12), font_color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.1), Inches(2.55), Inches(2.8), Inches(0.3),
                text=title, font_size=Pt(16), font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.1), Inches(2.9), Inches(2.8), Inches(0.9),
                text=desc, font_size=Pt(12), font_color=GRAY, alignment=PP_ALIGN.CENTER)

# A-round trigger box
trigger_box = add_rect(slide, Inches(0.6), Inches(4.2), Inches(12.1), Inches(2.2),
                       fill_color=BG_CARD, border_color=ACCENT2, border_width=Pt(1.5))

add_textbox(slide, Inches(0.9), Inches(4.35), Inches(5), Inches(0.35),
            text="A 轮触发条件（满足任意两项即可启动）", font_size=Pt(16), font_color=ACCENT2, bold=True)

triggers = [
    "□  MRR ≥ 50 万元，连续 3 个月",
    "□  MRR 月环比增长 ≥ 15%，连续 6 个月",
    "□  LTV/CAC ≥ 5:1  且  月流失率 ≤ 5%",
    "□  至少 4 张自有 GPU 稳定运营 3 个月，单卡毛利率 ≥ 40%",
]
for i, t in enumerate(triggers):
    add_textbox(slide, Inches(0.9), Inches(4.75 + i * 0.38), Inches(6), Inches(0.35),
                text=t, font_size=Pt(13), font_color=WHITE)

add_textbox(slide, Inches(8.0), Inches(4.35), Inches(4.5), Inches(1.0),
            text="A 轮目标\n1,500 ~ 2,500 万元\n用于首批 GPU 集群\n（20 ~ 40 张 H100/H20）",
            font_size=Pt(14), font_color=ACCENT, bold=True)


# ══════════════════════════════════════════════════════
# SLIDE 12 · 团队
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG_DARK)
add_slide_number(slide, 12)

add_section_title(slide, "团队——为什么这个人能做成")
# No subtitle needed

# Founder card
founder_card = add_rect(slide, Inches(0.6), Inches(1.5), Inches(6.0), Inches(4.5),
                        fill_color=BG_CARD, border_color=ACCENT, border_width=Pt(1))

add_textbox(slide, Inches(0.9), Inches(1.7), Inches(5.4), Inches(0.5),
            text="Eric  ·  创始人 & CEO", font_size=Pt(24), font_color=WHITE, bold=True)

add_textbox(slide, Inches(0.9), Inches(2.25), Inches(5.4), Inches(0.35),
            text="法律 × 金融 × 全栈开发  =  三重能力叠加", font_size=Pt(14), font_color=ACCENT)

capabilities = [
    ("合规壁垒", "算力交易涉及跨境 GPU 管制、数据安全法、\n个人信息保护法、ICP 许可证——\n这是我的本行。竞品复制不了。"),
    ("全栈工程能力", "FastAPI + React + PostgreSQL\n独立交付完整 SaaS 产品（档口通 CRM）\nMVP 不需要等 CTO，我就能写。"),
    ("商务谈判能力", "与云厂商谈代理协议、与投资人谈条款、\n与算力券管理部门对接——\n律师背景在这些场景中有天然优势。"),
]
for i, (cap_title, cap_desc) in enumerate(capabilities):
    y = Inches(2.8 + i * 1.1)
    add_textbox(slide, Inches(0.9), y, Inches(5.4), Inches(0.25),
                text=f"▸ {cap_title}", font_size=Pt(14), font_color=ACCENT2, bold=True)
    add_textbox(slide, Inches(0.9), y + Inches(0.3), Inches(5.4), Inches(0.7),
                text=cap_desc, font_size=Pt(12), font_color=GRAY)

# Right side: hiring plan
add_textbox(slide, Inches(7.2), Inches(1.5), Inches(5.5), Inches(0.35),
            text="种子轮拟招募", font_size=Pt(18), font_color=ACCENT2, bold=True)

hires = [
    ("技术合伙人 / CTO", "5% ~ 8% 股权 + 联合创始人 title\n3 年以上云计算/分布式系统经验\n熟悉 GPU 集群运维者优先"),
    ("全栈工程师 × 2", "月薪 2 万 + 期权 0.5%~1%\nReact + Python 后端\n通过开源项目贡献者筛选"),
    ("后端/基础设施 × 1", "月薪 2.2 万 + 期权 0.5%~1.5%\nKubernetes / Terraform / GPU 运维\n多云 SDK 开发经验"),
    ("社区运营 × 1", "兼职转全职\n内容营销 + 校园大使管理\nGPU 免费额度作为运营杠杆"),
]

for i, (role, desc) in enumerate(hires):
    y = Inches(2.1 + i * 1.1)
    add_rect(slide, Inches(7.2), y, Inches(5.5), Inches(0.9), fill_color=BG_CARD, border_color=BORDER, border_width=Pt(0.5))
    add_textbox(slide, Inches(7.45), y + Inches(0.08), Inches(5.0), Inches(0.25),
                text=role, font_size=Pt(13), font_color=ACCENT, bold=True)
    add_textbox(slide, Inches(7.45), y + Inches(0.38), Inches(5.0), Inches(0.5),
                text=desc, font_size=Pt(11), font_color=GRAY)

# Bottom
add_textbox(slide, Inches(0.6), Inches(6.3), Inches(12.1), Inches(0.5),
            text="竞品可以复制代码，但复制不了合规架构和商务谈判能力。",
            font_size=Pt(15), font_color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════
# SLIDE 13 · ASK
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG_DARK)
add_slide_number(slide, 13)

add_section_title(slide, "融资需求",
                  "种子轮 · 用 200 万验证 PMF，12 个月做到 A 轮就绪")

# Big ask
ask_box = add_rect(slide, Inches(0.6), Inches(1.5), Inches(12.1), Inches(1.0),
                   fill_color=BG_CARD, border_color=ACCENT2, border_width=Pt(2))

add_textbox(slide, Inches(1.0), Inches(1.6), Inches(11.3), Inches(0.7),
            text="200 万元  |  出让 10% ~ 15%  |  投前估值 1,133 ~ 1,800 万",
            font_size=Pt(26), font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Use of funds - bar chart style
funds = [
    ("研发人力（3人 × 12月）", "86.4 万", "43%", 8.6),
    ("GPU 采购押金", "40 万", "20%", 4.0),
    ("市场推广 & 社区运营", "24 万", "12%", 2.4),
    ("基础设施（云服务）", "12 万", "6%", 1.2),
    ("法务合规 & 公司注册", "5 万", "3%", 0.5),
    ("创始人生活费", "18 万", "9%", 1.8),
    ("机动预留", "14.6 万", "7%", 1.5),
]

colors_fund = [ACCENT, ACCENT2, RGBColor(0xFF, 0xD7, 0x00),
               GRAY, GRAY_DIM, GRAY, GRAY_DIM]

for i, (item, amount, pct, bar_w) in enumerate(funds):
    y = Inches(2.8 + i * 0.52)
    add_textbox(slide, Inches(0.8), y, Inches(3.5), Inches(0.35),
                text=item, font_size=Pt(12), font_color=WHITE)
    # Bar
    add_rect(slide, Inches(4.4), y + Inches(0.05), Inches(bar_w * 0.6), Inches(0.3),
             fill_color=colors_fund[i])
    add_textbox(slide, Inches(4.5 + bar_w * 0.6), y, Inches(1.8), Inches(0.35),
                text=amount, font_size=Pt(12), font_color=colors_fund[i], bold=True)
    add_textbox(slide, Inches(10.8), y, Inches(1.5), Inches(0.35),
                text=pct, font_size=Pt(12), font_color=WHITE)

# Bottom goal
goal_box = add_rect(slide, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.6),
                    fill_color=RGBColor(0x0A, 0x30, 0x20), border_color=ACCENT2, border_width=Pt(1))
add_textbox(slide, Inches(0.8), Inches(6.53), Inches(11.7), Inches(0.54),
            text="12 个月后目标：MRR 30 万  ·  盈亏平衡  ·  LTV/CAC > 3  ·  A 轮就绪",
            font_size=Pt(16), font_color=ACCENT2, bold=True, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════
# SLIDE 14 · 退出
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG_DARK)
add_slide_number(slide, 14)

add_section_title(slide, "退出路径与回报预期",
                  "算力赛道是 2026 年少有的退出路径清晰的领域")

acquirers = [
    ("阿里云", "补充“个人开发者”用户层\n大厂云触达不到的长尾", ACCENT),
    ("字节 · 火山引擎", "追赶阿里云/华为云\n需要一个现成的用户入口", ACCENT2),
    ("华为云", "昇腾国产芯片正在铺量\n需要算力平台作销售渠道", RGBColor(0xFF, 0xD7, 0x00)),
    ("AutoDL 等平台", "同类合并，覆盖全价格带\n消费级 + 企业级 GPU 互补", ACCENT3),
]

for i, (name, logic, color) in enumerate(acquirers):
    x = Inches(0.6 + i * 3.2)
    add_rect(slide, x, Inches(1.6), Inches(2.9), Inches(0.4), fill_color=color)
    add_textbox(slide, x + Inches(0.1), Inches(1.62), Inches(2.7), Inches(0.35),
                text=name, font_size=Pt(14), font_color=BG_DARK, bold=True)
    add_card(slide, x, Inches(2.1), Inches(2.9), Inches(1.4),
             title="", body_lines=[logic], body_size=Pt(12))

# Valuation math
add_textbox(slide, Inches(0.6), Inches(3.8), Inches(12.1), Inches(0.35),
            text="参考估值逻辑：算力服务平台通常按 ARR 的 5 ~ 10 倍估值", font_size=Pt(15), font_color=WHITE, bold=True)

# Calculation flow
calc_data = [
    "第三年基准方案 ARR  3,600 万",
    "↓  × 5~10 倍",
    "收购估值  1.8 亿 ~ 3.6 亿元",
    "↓  种子轮占股 10%~15%（稀释后约 7%~10%）",
    "种子轮 200 万 → 回报约 1,260 万 ~ 3,600 万",
    "→  回报倍数 6x ~ 18x",
]
for i, line in enumerate(calc_data):
    is_result = "→" in line
    add_textbox(slide, Inches(0.8), Inches(4.25 + i * 0.35), Inches(5.5), Inches(0.32),
                text=line, font_size=Pt(13), font_color=ACCENT2 if is_result else WHITE,
                bold=is_result)

# Alternative exits
add_textbox(slide, Inches(7.2), Inches(3.8), Inches(5.5), Inches(0.35),
            text="其他退出路径", font_size=Pt(15), font_color=ACCENT, bold=True)

alt_exits = [
    "路径 2：独立 IPO（5~8 年）\n科创板（““硬科技””定位）或港股\n参考：沐曦 GPU 已完成 IPO 辅导",
    "路径 3：传统 IDC / 运营商收购\n中国移动、万国数据、世纪互联\n正在从“卖机柜”向““卖算力””转型",
]
for i, txt in enumerate(alt_exits):
    add_card(slide, Inches(7.2), Inches(4.25 + i * 1.35), Inches(5.5), Inches(1.2),
             title="", body_lines=[txt], body_size=Pt(12))


# ══════════════════════════════════════════════════════
# SLIDE 15 · THANK YOU
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG_DARK)

# Top decorative line
add_line(slide, Inches(0), Inches(0), W, Pt(4), ACCENT)

add_textbox(slide, Inches(1.5), Inches(2.0), Inches(10.3), Inches(1.0),
            text="算 力 通", font_size=Pt(56), font_color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1.5), Inches(3.0), Inches(10.3), Inches(0.6),
            text="让每一个 AI 开发者，打开浏览器就能训练模型。",
            font_size=Pt(20), font_color=GRAY, alignment=PP_ALIGN.CENTER)

# Divider
add_line(slide, Inches(5.0), Inches(3.7), Inches(3.3), Pt(2), ACCENT)

add_textbox(slide, Inches(1.5), Inches(4.1), Inches(10.3), Inches(0.5),
            text="创始人：Eric", font_size=Pt(18), font_color=WHITE,
            alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1.5), Inches(4.6), Inches(10.3), Inches(0.5),
            text="邮箱：[待补充]    微信：[待补充]", font_size=Pt(14), font_color=GRAY,
            alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1.5), Inches(5.3), Inches(10.3), Inches(0.5),
            text="种子轮 · 200 万 · 10% ~ 15%", font_size=Pt(16), font_color=ACCENT2, bold=True,
            alignment=PP_ALIGN.CENTER)

# Bottom decorative line
add_line(slide, Inches(0), H - Pt(4), W, Pt(4), ACCENT)


# ══════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════
output_path = "/Users/apple/算力通/算力通_种子轮_Pitch_Deck.pptx"
prs.save(output_path)
print(f"✅ Pitch Deck 已生成: {output_path}")
print(f"   Slides: {len(prs.slides)}")
print("   Size: 16:9 widescreen")
