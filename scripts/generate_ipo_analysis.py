"""
算力通 · IPO 退出路径分析 PPT 生成器
对比并购退出，独立拆解 IPO 路径下的投资人回报
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Design Tokens ────────────────────────────────────
W = Inches(13.333)
H = Inches(7.5)

BG_DARK   = RGBColor(0x0A, 0x0E, 0x17)
BG_CARD   = RGBColor(0x12, 0x18, 0x24)
ACCENT    = RGBColor(0x00, 0xD4, 0xFF)
ACCENT2   = RGBColor(0x00, 0xFF, 0x88)
ACCENT3   = RGBColor(0xFF, 0x6B, 0x35)
ACCENT4   = RGBColor(0xFF, 0xD7, 0x00)
PURPLE    = RGBColor(0xA8, 0x55, 0xF7)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GRAY      = RGBColor(0x88, 0x93, 0xA0)
GRAY_DIM  = RGBColor(0x55, 0x5E, 0x6B)
BORDER    = RGBColor(0x1E, 0x28, 0x36)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank_layout = prs.slide_layouts[6]

# ── Helpers ──────────────────────────────────────────
def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fill_color=None, border_color=None, border_width=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.fill.solid()
        if border_width:
            shape.line.width = border_width
    return shape

def add_line(slide, l, t, w, h, color, lw=Pt(1)):
    return add_rect(slide, l, t, w, h, fill_color=color)

def add_textbox(slide, l, t, w, h, text="", font_size=Pt(14),
                font_color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = "PingFang SC"
    p.alignment = alignment
    p.line_spacing = Pt(18)
    return txBox

def add_ml_text(slide, l, t, w, h, lines, default_size=Pt(13),
                default_color=WHITE, default_bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(line, str):
            p.text = line
            p.font.size = default_size
            p.font.color.rgb = default_color
            p.font.bold = default_bold
        else:
            txt, sz, clr, bld = line[0], line[1] if len(line)>1 else default_size, \
                                line[2] if len(line)>2 else default_color, \
                                line[3] if len(line)>3 else default_bold
            p.text = txt
            p.font.size = sz
            p.font.color.rgb = clr
            p.font.bold = bld
        p.font.name = "PingFang SC"
        p.alignment = alignment
        p.line_spacing = Pt(18)
    return txBox

def add_section_title(slide, title, subtitle=""):
    add_line(slide, Inches(0), Inches(0), W, Pt(4), ACCENT)
    add_textbox(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.55),
                text=title, font_size=Pt(30), font_color=WHITE, bold=True)
    if subtitle:
        add_textbox(slide, Inches(0.6), Inches(0.95), Inches(12), Inches(0.35),
                    text=subtitle, font_size=Pt(14), font_color=GRAY)

def add_card(slide, l, t, w, h, title="", body_lines=None, title_color=ACCENT, body_size=Pt(11)):
    add_rect(slide, l, t, w, h, fill_color=BG_CARD, border_color=BORDER, border_width=Pt(0.5))
    if title:
        add_textbox(slide, l+Inches(0.15), t+Inches(0.1), w-Inches(0.3), Inches(0.28),
                    text=title, font_size=Pt(13), font_color=title_color, bold=True)
    if body_lines:
        add_ml_text(slide, l+Inches(0.15), t+Inches(0.42), w-Inches(0.3), h-Inches(0.5),
                    lines=body_lines, default_size=body_size, default_color=GRAY)


# ══════════════════════════════════════════════════════
# SLIDE 1 · COVER
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG_DARK)
add_line(slide, Inches(0), Inches(0), W, Pt(6), PURPLE)
add_textbox(slide, Inches(1), Inches(1.8), Inches(11.3), Inches(0.9),
            text="算力通 · IPO 退出路径分析", font_size=Pt(44), font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_line(slide, Inches(4.5), Inches(2.9), Inches(4.3), Pt(2), PURPLE)
add_ml_text(slide, Inches(1.5), Inches(3.3), Inches(10.3), Inches(1.5),
            lines=[
                ("科创板 / 港股  ·  7-8 年路径  ·  回报率测算", Pt(20), GRAY, False),
                ("", Pt(10), GRAY, False),
                ("与 3 年并购退出并行对照  |  种子轮 · 200 万 · 10%~15%", Pt(14), PURPLE, False),
            ], alignment=PP_ALIGN.CENTER)
add_line(slide, Inches(0), H-Pt(6), W, Pt(6), PURPLE)


# ══════════════════════════════════════════════════════
# SLIDE 2 · IPO vs M&A 本质差异
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG_DARK)
add_section_title(slide, "IPO vs 并购：两条退出路径的本质差异",
                  "同一个项目，两种退出哲学 — IRR 与 MOIC 的取舍")

# Comparison table
headers = ["维度", "并购（3-5年）", "IPO（7-8年）", "差异"]
rows = [
    ("退出方式", "一次性现金/换股", "公开市场逐步减持", "流动性不同"),
    ("估值倍数", "ARR 5-10x", "ARR 8-15x", "IPO 高 60%+"),
    ("锁定期", "无", "6-12 个月", "IPO 有流动性限制"),
    ("融资轮次", "种子 + A 轮", "种子 + A + B + 可能 C", "IPO 多 2-3 轮"),
    ("累计稀释", "~24%", "~55-60%", "IPO 多稀释 2.5x"),
    ("种子最终持股", "~9.5%", "~4.9%", "并购是 IPO 的 1.9 倍"),
    ("MOIC（基准）", "12.9x", "25.9x", "IPO 高 2x"),
    ("IRR（基准）", "134%", "59%", "并购高 2.3x"),
    ("确定性", "较高", "中低", "并购更可控"),
]

col_w = [Inches(2.0), Inches(3.3), Inches(3.5), Inches(3.5)]
for r, row in enumerate([headers] + rows):
    y = Inches(1.5 + r * 0.52)
    is_header = r == 0
    for c, cell in enumerate(row):
        x = Inches(0.6) + sum(cw for cw in col_w[:c])
        if is_header:
            add_rect(slide, x, y, col_w[c], Inches(0.45), fill_color=PURPLE)
            add_textbox(slide, x, y+Inches(0.06), col_w[c], Inches(0.33),
                        text=cell, font_size=Pt(11), font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        else:
            bg = BG_DARK if r % 2 == 0 else BG_CARD
            is_highlight = r in [7, 8]  # MOIC and IRR rows
            txt_clr = ACCENT2 if is_highlight else WHITE
            bld = is_highlight
            add_rect(slide, x, y, col_w[c], Inches(0.45), fill_color=bg, border_color=BORDER, border_width=Pt(0.3))
            add_textbox(slide, x+Inches(0.05), y+Inches(0.06), col_w[c]-Inches(0.1), Inches(0.33),
                        text=cell, font_size=Pt(11), font_color=txt_clr, bold=bld, alignment=PP_ALIGN.CENTER)

# Bottom insight
add_rect(slide, Inches(0.6), Inches(6.3), Inches(12.1), Inches(0.8), fill_color=BG_CARD, border_color=PURPLE)
add_ml_text(slide, Inches(0.9), Inches(6.35), Inches(11.5), Inches(0.7),
            lines=[
                ("核心悖论：并购 IRR 高但倍数有限，IPO 倍数高但 IRR 被时间稀释。", Pt(15), ACCENT2, True),
                ("种子投资人选哪个，取决于基金存续期和组合策略 — 没有绝对好坏，只有是否匹配。", Pt(12), GRAY, False),
            ])


# ══════════════════════════════════════════════════════
# SLIDE 3 · IPO 时间线
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG_DARK)
add_section_title(slide, "IPO 路径时间线：5 年太紧，7-8 年更现实",
                  "科创板/港股上市门槛 + 业务增速衰减 → 大概率在 Y7-Y8 挂牌")

# Timeline visualization
milestones = [
    ("2026.Q3", "种子轮\n200万", PURPLE, Inches(2.2)),
    ("2027.Q4", "A 轮\n2,000万", ACCENT, Inches(2.2)),
    ("2029.Q1", "B 轮\n7,500万", ACCENT2, Inches(2.2)),
    ("2031", "C 轮(可能)\n1-2亿", ACCENT4, Inches(2.2)),
    ("2032-33", "IPO申报\n科创板/港股", ACCENT3, Inches(2.2)),
    ("2033-34", "挂牌交易\n市值12-25亿", PURPLE, Inches(2.2)),
]

for i, (time, event, color, h) in enumerate(milestones):
    x = Inches(0.4 + i * 2.15)
    # connector line
    if i < 5:
        add_line(slide, x + Inches(1.9), Inches(2.7), Inches(0.35), Pt(3), color)
    # node
    add_rect(slide, x + Inches(0.7), Inches(2.3), Inches(0.25), Inches(0.25), fill_color=color)
    # box
    add_rect(slide, x, Inches(2.7), Inches(1.95), h, fill_color=color)
    add_textbox(slide, x+Inches(0.05), Inches(2.75), Inches(1.85), Inches(0.3),
                text=time, font_size=Pt(11), font_color=BG_DARK, bold=True, alignment=PP_ALIGN.CENTER)
    add_ml_text(slide, x+Inches(0.05), Inches(3.1), Inches(1.85), h-Inches(0.5),
                lines=[(event, Pt(11), BG_DARK, False)], alignment=PP_ALIGN.CENTER)

# 5-year vs 7-year comparison
add_textbox(slide, Inches(0.6), Inches(5.3), Inches(5.8), Inches(0.3),
            text="❌  为什么 5 年 IPO 极难？", font_size=Pt(16), font_color=ACCENT3, bold=True)
add_ml_text(slide, Inches(0.6), Inches(5.7), Inches(5.8), Inches(1.5),
            lines=[
                ("科创板「软件企业」隐含门槛：", Pt(12), WHITE, True),
                ("· 年营收 3-5 亿 + 连续盈利", Pt(11), GRAY, False),
                ("· 或年营收 5 亿+ 且高增长（>30%）", Pt(11), GRAY, False),
                ("· 5 年内从零跑到这个量级 → 必须连续翻倍", Pt(11), ACCENT3, False),
                ("→ 概率偏低，不适合作为融资假设", Pt(11), ACCENT3, True),
            ])

add_textbox(slide, Inches(6.8), Inches(5.3), Inches(5.8), Inches(0.3),
            text="✅  为什么 7-8 年更现实？", font_size=Pt(16), font_color=ACCENT2, bold=True)
add_ml_text(slide, Inches(6.8), Inches(5.7), Inches(5.8), Inches(1.5),
            lines=[
                ("Y7 ARR ~1.18 亿 + 高毛利率 + 盈利", Pt(12), WHITE, True),
                ("· 科创板「关键软件技术」定位可行", Pt(11), GRAY, False),
                ("· 港股 SaaS 门槛更低、审核更快", Pt(11), GRAY, False),
                ("· 增速 20%+ 仍属于高增长公司", Pt(11), ACCENT2, False),
                ("→ 这是建模的基准假设", Pt(11), ACCENT2, True),
            ])


# ══════════════════════════════════════════════════════
# SLIDE 4 · ARR 增速衰减曲线
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG_DARK)
add_section_title(slide, "ARR 超 Y3 之后：增速自然衰减曲线",
                  "从 +150% 到 +15% — 大数法则下 Y4-Y8 外推逻辑")

arr_data = [
    ("Y1", "360万", "—", "BP 基准", ACCENT),
    ("Y2", "1,440万", "+300%", "BP 基准 · 从0到1", ACCENT),
    ("Y3", "3,600万", "+150%", "BP 基准 · 双轨初显", ACCENT2),
    ("Y4", "5,400万", "+50%", "增速首次跌破100%", ACCENT),
    ("Y5", "7,560万", "+40%", "存量增长为主", ACCENT),
    ("Y6", "9,830万", "+30%", "接近SaaS成熟增速", ACCENT),
    ("Y7", "1.18亿", "+20%", "IPO申报窗口 ★", ACCENT2),
    ("Y8", "1.36亿", "+15%", "挂牌当年", PURPLE),
]

for i, (yr, arr, growth, note, color) in enumerate(arr_data):
    y = Inches(1.5 + i * 0.62)
    is_divider = i == 3
    if is_divider:
        add_line(slide, Inches(0.6), y - Inches(0.05), Inches(12.1), Pt(1), GRAY_DIM)
    add_rect(slide, Inches(0.6), y, Inches(12.1), Inches(0.55), fill_color=BG_CARD if i%2==0 else BG_DARK, border_color=BORDER, border_width=Pt(0.3))
    # Year badge
    add_rect(slide, Inches(0.7), y+Inches(0.08), Inches(0.9), Inches(0.38), fill_color=color)
    add_textbox(slide, Inches(0.7), y+Inches(0.1), Inches(0.9), Inches(0.34),
                text=yr, font_size=Pt(13), font_color=BG_DARK, bold=True, alignment=PP_ALIGN.CENTER)
    # ARR
    add_textbox(slide, Inches(1.8), y+Inches(0.1), Inches(2.0), Inches(0.34),
                text=arr, font_size=Pt(16), font_color=WHITE, bold=True)
    # Growth
    txt_clr = ACCENT2 if "+" in growth else GRAY
    add_textbox(slide, Inches(3.9), y+Inches(0.1), Inches(1.5), Inches(0.34),
                text=growth, font_size=Pt(13), font_color=txt_clr, bold=True)
    # Note
    add_textbox(slide, Inches(5.5), y+Inches(0.1), Inches(7.0), Inches(0.34),
                text=note, font_size=Pt(11), font_color=GRAY)

# CAGR callout
add_rect(slide, Inches(0.6), Inches(6.6), Inches(12.1), Inches(0.55), fill_color=BG_CARD, border_color=ACCENT2)
add_textbox(slide, Inches(0.9), Inches(6.65), Inches(11.5), Inches(0.45),
            text="Y3→Y7 CAGR ≈ 35%  |  显著低于智能算力市场 46% CAGR  |  保守合理  |  Y3 基准 3,600 万是可信的 IPO 起跑线",
            font_size=Pt(13), font_color=WHITE, bold=True)


# ══════════════════════════════════════════════════════
# SLIDE 5 · 完整稀释路径
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG_DARK)
add_section_title(slide, "稀释全路径：12.5% 走到 IPO 还剩 4.9%",
                  "Seed → A → B → C → IPO 发行  |  累计稀释 ~61%，但饼大了 60+ 倍")

dilution = [
    ("M0", "种子轮 close", "12.50%", "—", "—", "投后 1,600万", PURPLE),
    ("M18", "期权池 15%→20% + A轮", "9.35%", "-7.9%", "-18.8%", "投后 1.06亿", ACCENT),
    ("M36", "期权池 refresh + B轮", "7.05%", "-8.0%", "-18.0%", "投后 4.35亿", ACCENT2),
    ("M60", "期权池 refresh + C轮(可能)", "5.74%", "-7.1%", "-12.5%", "投后 ~13.5亿", ACCENT4),
    ("M84", "IPO 发行 15% 新股", "4.88%", "—", "-15.0%", "市值 10.62亿", PURPLE),
]

col_w2 = [Inches(0.7), Inches(3.0), Inches(1.4), Inches(1.3), Inches(1.3), Inches(2.0), Inches(1.8)]
headers2 = ["", "事件", "持股", "扩池稀释", "新股稀释", "公司估值", ""]

for r, row in enumerate([headers2] + dilution):
    y = Inches(1.5 + r * 0.65)
    is_header = r == 0
    if is_header:
        cells = row
        for c, cell in enumerate(cells):
            x = Inches(0.6) + sum(cw for cw in col_w2[:c])
            add_rect(slide, x, y, col_w2[c], Inches(0.42), fill_color=PURPLE)
            add_textbox(slide, x, y+Inches(0.05), col_w2[c], Inches(0.32),
                        text=cell, font_size=Pt(10), font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    else:
        time, event, stake, pool_dil, new_dil, val, color = row
        cells = [time, event, stake, pool_dil, new_dil, val, ""]
        for c, cell in enumerate(cells):
            x = Inches(0.6) + sum(cw for cw in col_w2[:c])
            bg = BG_DARK if r % 2 == 0 else BG_CARD
            if c == 2:
                clr = ACCENT2
                bld = True
                fs = Pt(15)
            elif c == 5:
                clr = WHITE
                bld = True
                fs = Pt(11)
            else:
                clr = WHITE if c in [0,1] else GRAY
                bld = False
                fs = Pt(11)
            add_rect(slide, x, y, col_w2[c], Inches(0.55), fill_color=bg, border_color=BORDER, border_width=Pt(0.3))
            add_textbox(slide, x+Inches(0.05), y+Inches(0.08), col_w2[c]-Inches(0.1), Inches(0.38),
                        text=cell, font_size=fs, font_color=clr, bold=bld, alignment=PP_ALIGN.CENTER)

# Summary
add_rect(slide, Inches(0.6), Inches(5.9), Inches(5.8), Inches(1.2), fill_color=BG_CARD, border_color=ACCENT2)
add_ml_text(slide, Inches(0.9), Inches(5.95), Inches(5.2), Inches(1.1),
            lines=[
                ("股权变小了，但饼大了 60+ 倍", Pt(18), ACCENT2, True),
                ("种子轮 12.5% → IPO 时 4.88%", Pt(14), GRAY, False),
                ("但公司从 1,600 万 → 10.62 亿", Pt(14), GRAY, False),
                ("每一轮「被稀释」= 估值大幅跃升", Pt(12), GRAY_DIM, False),
            ])

add_rect(slide, Inches(6.8), Inches(5.9), Inches(5.9), Inches(1.2), fill_color=BG_CARD, border_color=ACCENT3)
add_ml_text(slide, Inches(7.1), Inches(5.95), Inches(5.3), Inches(1.1),
            lines=[
                ("对比并购路径：", Pt(16), ACCENT3, True),
                ("并购稀释 ~24% → 最终持股 9.55%", Pt(13), GRAY, False),
                ("IPO 稀释 ~61% → 最终持股 4.88%", Pt(13), GRAY, False),
                ("但 IPO 估值是并购的 3.9 倍", Pt(12), ACCENT2, True),
            ])


# ══════════════════════════════════════════════════════
# SLIDE 6 · 可比 IPO 估值锚
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG_DARK)
add_section_title(slide, "IPO 估值锚：算力平台值多少倍 ARR？",
                  "对标 CoreWeave / Paperspace / Lambda  |  科创板/港股折价调整")

comps = [
    ("公司", "事件/年份", "市值/估值", "ARR倍数", "毛利率", "增速", "相关性"),
    ("CoreWeave", "IPO 2025", "~230亿美元", "~12x", "~75%", "+100%+", "GPU云 · 高度相关"),
    ("Lambda Labs", "融资 2025", "~15亿美元", "~10x", "~60%", "+80%", "GPU租赁 · 直接参考"),
    ("Paperspace", "被收购 2023", "1.11亿美元", "~8x", "~55%", "+30%", "GPU平台 · 高度相关"),
    ("UCloud", "科创板 2020", "~150亿", "~8x", "~15%", "+30%", "云计算 · 低毛利参照"),
    ("金山云", "港股 2022", "~80亿港元", "~3x", "~5%", "+10%", "纯IaaS · 不可比"),
]

cw3 = [Inches(1.5), Inches(1.5), Inches(1.8), Inches(1.4), Inches(1.1), Inches(1.2), Inches(3.5)]
for r, row in enumerate(comps):
    y = Inches(1.5 + r * 0.52)
    is_header = r == 0
    for c, cell in enumerate(row):
        x = Inches(0.6) + sum(cw for cw in cw3[:c])
        if is_header:
            add_rect(slide, x, y, cw3[c], Inches(0.42), fill_color=PURPLE)
            add_textbox(slide, x, y+Inches(0.04), cw3[c], Inches(0.34),
                        text=cell, font_size=Pt(10), font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        else:
            bg = BG_DARK if r % 2 == 0 else BG_CARD
            txt_clr = ACCENT2 if "~12x" in cell or "~10x" in cell else WHITE
            add_rect(slide, x, y, cw3[c], Inches(0.45), fill_color=bg, border_color=BORDER, border_width=Pt(0.3))
            add_textbox(slide, x+Inches(0.05), y+Inches(0.06), cw3[c]-Inches(0.1), Inches(0.33),
                        text=cell, font_size=Pt(10), font_color=txt_clr, alignment=PP_ALIGN.CENTER)

# Analysis box
add_rect(slide, Inches(0.6), Inches(4.8), Inches(6.0), Inches(2.3), fill_color=BG_CARD, border_color=ACCENT2)
add_textbox(slide, Inches(0.9), Inches(4.9), Inches(5.4), Inches(0.3),
            text="算力通对标分析", font_size=Pt(16), font_color=ACCENT2, bold=True)
add_ml_text(slide, Inches(0.9), Inches(5.3), Inches(5.4), Inches(1.7),
            lines=[
                ("❌ 不对标金山云（低毛利 IaaS, 3x）", Pt(12), ACCENT3, False),
                ("  算力通毛利率 20%→42%→55%，远超纯 IaaS", Pt(11), GRAY, False),
                ("", Pt(6), GRAY, False),
                ("❌ 不直接对标 CoreWeave（12x）", Pt(12), ACCENT3, False),
                ("  CoreWeave 自有万卡集群，重资产高壁垒", Pt(11), GRAY, False),
                ("", Pt(6), GRAY, False),
                ("✅ 合理区间：Paperspace 8x ~ Lambda 10x", Pt(13), ACCENT2, True),
                ("  轻资产 GPU 平台，毛利率 40-55%", Pt(11), GRAY, False),
            ])

add_rect(slide, Inches(6.9), Inches(4.8), Inches(5.8), Inches(2.3), fill_color=BG_CARD, border_color=PURPLE)
add_textbox(slide, Inches(7.2), Inches(4.9), Inches(5.2), Inches(0.3),
            text="A股/港股流动性折价", font_size=Pt(16), font_color=PURPLE, bold=True)
add_ml_text(slide, Inches(7.2), Inches(5.3), Inches(5.2), Inches(1.7),
            lines=[
                ("中国公开市场 vs 美股的估值折价：", Pt(12), WHITE, True),
                ("· 科创板 SaaS：通常 6-10x ARR", Pt(11), GRAY, False),
                ("· 港股 SaaS：通常 5-8x ARR", Pt(11), GRAY, False),
                ("· 比美股同类低 10-25%", Pt(11), GRAY, False),
                ("", Pt(6), GRAY, False),
                ("算力通 IPO 估值倍数：", Pt(13), ACCENT2, True),
                ("8-10x ARR（取 9x 中值）", Pt(18), PURPLE, True),
            ])


# ══════════════════════════════════════════════════════
# SLIDE 7 · 全情景回报矩阵
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG_DARK)
add_section_title(slide, "IPO 退出：全情景回报矩阵",
                  "种子轮 200 万 → 稀释至 4.88% → 7 年退出")

add_textbox(slide, Inches(0.6), Inches(1.4), Inches(12.1), Inches(0.3),
            text="IPO 市值 = ARR × 倍数  |  种子回报 = IPO 市值 × 4.88%",
            font_size=Pt(13), font_color=GRAY, bold=True)

scenarios = [
    ("情景", "ARR(亿)", "倍数", "IPO市值", "种子回报", "MOIC", "7年IRR"),
    ("保守 · Y8", "0.98", "7x", "6.86亿", "3,348万", "16.7x", "50%"),
    ("保守 · Y7", "0.85", "8x", "6.80亿", "3,318万", "16.6x", "49%"),
    ("基准 · Y8", "1.36", "8x", "10.88亿", "5,309万", "26.5x", "57%"),
    ("基准 · Y7 ★", "1.18", "9x", "10.62亿", "5,183万", "25.9x", "59%"),
    ("基准 · Y7", "1.18", "10x", "11.80亿", "5,758万", "28.8x", "61%"),
    ("乐观 · Y7", "2.10", "10x", "21.00亿", "10,248万", "51.2x", "76%"),
    ("乐观 · Y6(无C轮)", "1.60", "11x", "17.60亿", "10,102万", "50.5x", "89%"),
]

scol_w = [Inches(1.8), Inches(1.3), Inches(1.1), Inches(1.8), Inches(2.0), Inches(1.3), Inches(1.6)]
for r, row in enumerate(scenarios):
    y = Inches(1.85 + r * 0.52)
    is_h = r == 0
    for c, cell in enumerate(row):
        x = Inches(0.6) + sum(cw for cw in scol_w[:c])
        if is_h:
            add_rect(slide, x, y, scol_w[c], Inches(0.42), fill_color=PURPLE)
            add_textbox(slide, x, y+Inches(0.04), scol_w[c], Inches(0.34),
                        text=cell, font_size=Pt(10), font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        else:
            is_star = "★" in cell
            bg = RGBColor(0x15, 0x08, 0x28) if is_star else (BG_DARK if r%2==0 else BG_CARD)
            txt_clr = ACCENT2 if (is_star or ("x" in cell and "亿" not in cell and "x" in cell)) else WHITE
            bld = is_star or (c >= 5)
            add_rect(slide, x, y, scol_w[c], Inches(0.45), fill_color=bg, border_color=BORDER, border_width=Pt(0.3))
            add_textbox(slide, x, y+Inches(0.06), scol_w[c], Inches(0.33),
                        text=cell, font_size=Pt(12) if is_star else Pt(11), font_color=txt_clr, bold=bld, alignment=PP_ALIGN.CENTER)

# Key takeaway
add_rect(slide, Inches(0.6), Inches(5.75), Inches(12.1), Inches(1.4), fill_color=BG_CARD, border_color=PURPLE)
add_ml_text(slide, Inches(0.9), Inches(5.85), Inches(11.5), Inches(1.2),
            lines=[
                ("基准方案（最可能）：Y7 IPO · ARR 1.18 亿 · 9x → 市值 10.62 亿 → 种子回报 5,183 万 → 25.9x MOIC", Pt(16), PURPLE, True),
                ("", Pt(6), GRAY, False),
                ("保守方案底线：Y8 以 7x 上市 → 仍有 16.7x MOIC → 远超 VC 10x 目标", Pt(13), WHITE, False),
                ("乐观方案上限：Y7 以 10x 上市 → 51.2x MOIC → 200 万变 1 亿，基金回报级项目", Pt(13), ACCENT2, False),
            ])


# ══════════════════════════════════════════════════════
# SLIDE 8 · 并购 vs IPO 并行对照
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG_DARK)
add_section_title(slide, "并购 vs IPO：同一项目，两条回报曲线",
                  "投资人怎么选？取决于基金存续期、风险偏好、组合策略")

# Side by side
# Left: M&A
add_rect(slide, Inches(0.4), Inches(1.4), Inches(5.8), Inches(0.55), fill_color=ACCENT)
add_textbox(slide, Inches(0.5), Inches(1.43), Inches(5.6), Inches(0.5),
            text="路径一：并购退出（3年）", font_size=Pt(20), font_color=BG_DARK, bold=True, alignment=PP_ALIGN.CENTER)

ma_metrics = [
    ("种子持股", "9.55%"),
    ("退出估值", "2.70 亿 (ARR 3,600万 × 7.5x)"),
    ("种子回报", "2,579 万"),
    ("MOIC", "12.9x"),
    ("IRR (3年)", "134%"),
    ("确定性", "⭐⭐⭐⭐  较高"),
]
for i, (label, value) in enumerate(ma_metrics):
    y = Inches(2.1 + i * 0.5)
    add_textbox(slide, Inches(0.7), y, Inches(2.0), Inches(0.35),
                text=label, font_size=Pt(14), font_color=ACCENT, bold=True)
    add_textbox(slide, Inches(2.8), y, Inches(3.2), Inches(0.35),
                text=value, font_size=Pt(14), font_color=WHITE, bold=(i>=3))

# Right: IPO
add_rect(slide, Inches(6.8), Inches(1.4), Inches(5.8), Inches(0.55), fill_color=PURPLE)
add_textbox(slide, Inches(6.9), Inches(1.43), Inches(5.6), Inches(0.5),
            text="路径二：IPO 退出（7年）", font_size=Pt(20), font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

ipo_metrics = [
    ("种子持股", "4.88%"),
    ("退出估值", "10.62 亿 (ARR 1.18亿 × 9x)"),
    ("种子回报", "5,183 万"),
    ("MOIC", "25.9x"),
    ("IRR (7年)", "59%"),
    ("确定性", "⭐⭐⭐  中低"),
]
for i, (label, value) in enumerate(ipo_metrics):
    y = Inches(2.1 + i * 0.5)
    add_textbox(slide, Inches(7.1), y, Inches(2.0), Inches(0.35),
                text=label, font_size=Pt(14), font_color=PURPLE, bold=True)
    add_textbox(slide, Inches(9.2), y, Inches(3.2), Inches(0.35),
                text=value, font_size=Pt(14), font_color=WHITE, bold=(i>=3))

# Arrows highlighting key differences
add_line(slide, Inches(0.5), Inches(5.0), Inches(12.3), Pt(1), GRAY_DIM)

add_ml_text(slide, Inches(0.6), Inches(5.15), Inches(5.5), Inches(2.0),
            lines=[
                ("并购的优势", Pt(18), ACCENT, True),
                ("", Pt(6), GRAY, False),
                ("✅ IRR 134% — 钱回来得快", Pt(13), ACCENT2, True),
                ("✅ 确定性高 — 有明确买家", Pt(13), WHITE, False),
                ("✅ 锁定期无 — 一次性变现", Pt(13), WHITE, False),
                ("", Pt(6), GRAY, False),
                ("⚠️ 倍数受限 — 把未来增长卖给了收购方", Pt(12), GRAY, False),
            ])

add_ml_text(slide, Inches(6.8), Inches(5.15), Inches(5.5), Inches(2.0),
            lines=[
                ("IPO 的优势", Pt(18), PURPLE, True),
                ("", Pt(6), GRAY, False),
                ("✅ MOIC 25.9x — 绝对回报更高", Pt(13), ACCENT2, True),
                ("✅ 7 年复合增长全归现有股东", Pt(13), WHITE, False),
                ("✅ 公开市场流动性 + 品牌溢价", Pt(13), WHITE, False),
                ("", Pt(6), GRAY, False),
                ("⚠️ IRR 低 — 时间稀释回报", Pt(12), GRAY, False),
                ("⚠️ 确定性低 — 需满足上市条件", Pt(12), GRAY, False),
            ])


# ══════════════════════════════════════════════════════
# SLIDE 9 · 投资人决策框架
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG_DARK)
add_section_title(slide, "投资人决策框架：你的基金适合哪条路？",
                  "没有绝对好坏 — 只有是否匹配基金策略")

decisions = [
    ("存续期还剩 3-4 年", "→ 并购", "必须在基金到期前退出", ACCENT),
    ("存续期还剩 7-8 年", "→ IPO", "多等 4 年换 2x 额外回报", PURPLE),
    ("追求高 IRR (>100%)", "→ 并购", "3 年 134% IRR 是顶级表现", ACCENT),
    ("追求高 MOIC (>20x)", "→ IPO", "25.9x 是基金回报级项目", PURPLE),
    ("风险厌恶型 LP", "→ 并购", "确定性远高于 IPO", ACCENT),
    ("愿意承担公开市场风险", "→ IPO", "流动性 + 上行空间更大", PURPLE),
    ("需要 DPI 证明", "→ 并购", "3 年现金回笼，LP 看得见", ACCENT),
    ("追求 TVPI 最大化", "→ IPO", "长线持有享受复利效应", PURPLE),
]

for i, (condition, choice, reason, color) in enumerate(decisions):
    y = Inches(1.5 + i * 0.58)
    add_rect(slide, Inches(0.6), y, Inches(5.5), Inches(0.5), fill_color=BG_CARD, border_color=color, border_width=Pt(1))
    add_textbox(slide, Inches(0.8), y+Inches(0.06), Inches(5.1), Inches(0.38),
                text=condition, font_size=Pt(13), font_color=WHITE, bold=True)
    add_rect(slide, Inches(6.3), y, Inches(1.3), Inches(0.5), fill_color=color)
    add_textbox(slide, Inches(6.3), y+Inches(0.06), Inches(1.3), Inches(0.38),
                text=choice, font_size=Pt(14), font_color=BG_DARK, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(7.8), y+Inches(0.06), Inches(5.0), Inches(0.38),
                text=reason, font_size=Pt(12), font_color=GRAY)

# Bottom summary
add_rect(slide, Inches(0.6), Inches(6.3), Inches(12.1), Inches(0.8), fill_color=BG_CARD, border_color=PURPLE, border_width=Pt(1.5))
add_ml_text(slide, Inches(0.9), Inches(6.35), Inches(11.5), Inches(0.7),
            lines=[
                ("两条路径不是互斥的。最优策略：做到 Y3-Y4 时接并购邀约 → 出价合理就卖 → 没有就融 B 轮继续做 IPO。", Pt(15), ACCENT2, True),
                ("种子投资人无论走哪条路，基准回报都在 12x 以上 — 这在种子轮中属于前四分之一分位。", Pt(12), GRAY, False),
            ])


# ══════════════════════════════════════════════════════
# SLIDE 10 · 综合期望回报
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG_DARK)
add_section_title(slide, "综合期望回报：加权后的种子轮预期",
                  "并购(60-70%概率) + IPO(15-25%) + 其他(10-15%) → 加权 MOIC")

# Probability-weighted return
add_textbox(slide, Inches(0.6), Inches(1.5), Inches(12.1), Inches(0.35),
            text="三条退出路径的加权期望", font_size=Pt(18), font_color=WHITE, bold=True)

weighted = [
    ("退出路径", "概率", "MOIC", "种子回报", "加权贡献"),
    ("路径一 · 并购（3年）", "60-70%", "12.9x", "2,579万", "~1,680万"),
    ("路径二 · IPO（7年）", "15-25%", "25.9x", "5,183万", "~1,037万"),
    ("路径三 · 其他/失败", "10-15%", "0-3x", "0-600万", "~60万"),
    ("", "", "", "", ""),
    ("加权合计", "100%", "~14-17x", "~2,800万", "—"),
]

ww = [Inches(2.5), Inches(1.5), Inches(1.5), Inches(2.5), Inches(2.5)]
for r, row in enumerate(weighted):
    y = Inches(2.0 + r * 0.52)
    is_header = r == 0
    is_total = r == 5
    for c, cell in enumerate(row):
        x = Inches(0.6) + sum(cw for cw in ww[:c])
        if is_header:
            add_rect(slide, x, y, ww[c], Inches(0.45), fill_color=PURPLE)
            add_textbox(slide, x, y+Inches(0.06), ww[c], Inches(0.33),
                        text=cell, font_size=Pt(12), font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        elif is_total:
            add_rect(slide, x, y, ww[c], Inches(0.5), fill_color=ACCENT2)
            add_textbox(slide, x, y+Inches(0.06), ww[c], Inches(0.38),
                        text=cell, font_size=Pt(15), font_color=BG_DARK, bold=True, alignment=PP_ALIGN.CENTER)
        else:
            bg = BG_DARK if r % 2 == 0 else BG_CARD
            add_rect(slide, x, y, ww[c], Inches(0.45), fill_color=bg, border_color=BORDER, border_width=Pt(0.3))
            add_textbox(slide, x+Inches(0.05), y+Inches(0.06), ww[c]-Inches(0.1), Inches(0.33),
                        text=cell, font_size=Pt(12), font_color=ACCENT2 if c==2 else WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Context
add_rect(slide, Inches(0.6), Inches(4.9), Inches(5.8), Inches(2.2), fill_color=BG_CARD, border_color=ACCENT2)
add_textbox(slide, Inches(0.9), Inches(5.0), Inches(5.2), Inches(0.3),
            text="这在 VC 中是什么水平？", font_size=Pt(16), font_color=ACCENT2, bold=True)
add_ml_text(slide, Inches(0.9), Inches(5.4), Inches(5.2), Inches(1.5),
            lines=[
                ("种子轮 MOIC 分位数：", Pt(13), WHITE, True),
                ("· 前 10%：>20x", Pt(12), GRAY, False),
                ("· 前 25%：>10x", Pt(12), GRAY, False),
                ("· 中位数：~3x", Pt(12), GRAY, False),
                ("· 后 25%：<1x（亏损）", Pt(12), GRAY, False),
                ("", Pt(6), GRAY, False),
                ("算力通加权 MOIC ~14-17x", Pt(15), ACCENT2, True),
                ("属于前 15-20% 分位", Pt(13), ACCENT2, True),
            ])

add_rect(slide, Inches(6.8), Inches(4.9), Inches(5.9), Inches(2.2), fill_color=BG_CARD, border_color=PURPLE)
add_textbox(slide, Inches(7.1), Inches(5.0), Inches(5.3), Inches(0.3),
            text="为什么算力通的期望回报偏高？", font_size=Pt(16), font_color=PURPLE, bold=True)
add_ml_text(slide, Inches(7.1), Inches(5.4), Inches(5.3), Inches(1.5),
            lines=[
                ("三个结构性优势：", Pt(13), WHITE, True),
                ("", Pt(6), GRAY, False),
                ("1. 下行保护极强", Pt(13), ACCENT2, True),
                ("   ARR 280万即回本，本金损失概率<5%", Pt(11), GRAY, False),
                ("2. 退出路径清晰", Pt(13), ACCENT2, True),
                ("   有明确收购方 + 公开可比交易", Pt(11), GRAY, False),
                ("3. 资本效率高", Pt(13), ACCENT2, True),
                ("   LTV/CAC=16:1，不需烧钱换增长", Pt(11), GRAY, False),
            ])


# ══════════════════════════════════════════════════════
# SLIDE 11 · END PAGE
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG_DARK)
add_line(slide, Inches(0), Inches(0), W, Pt(6), PURPLE)
add_textbox(slide, Inches(1.5), Inches(1.8), Inches(10.3), Inches(0.9),
            text="算力通 · IPO 退出路径分析", font_size=Pt(42), font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_line(slide, Inches(4.5), Inches(2.9), Inches(4.3), Pt(2), PURPLE)
add_ml_text(slide, Inches(1.5), Inches(3.3), Inches(10.3), Inches(1.8),
            lines=[
                ("并购 12.9x MOIC / 134% IRR  ←→  IPO 25.9x MOIC / 59% IRR", Pt(18), GRAY, False),
                ("", Pt(10), GRAY, False),
                ("本报告基于商业计划书 V3 数据与公开市场可比分析", Pt(14), GRAY_DIM, False),
                ("不构成投资建议 · 仅供种子轮投资人审阅", Pt(14), GRAY_DIM, False),
                ("", Pt(10), GRAY, False),
                ("算力通 · 让算力像自来水一样，扭开即用", Pt(18), PURPLE, True),
            ], alignment=PP_ALIGN.CENTER)
add_line(slide, Inches(5.5), Inches(5.5), Inches(2.3), Pt(2), PURPLE)
add_textbox(slide, Inches(1.5), Inches(5.8), Inches(10.3), Inches(0.5),
            text="创始人：Eric  |  种子轮 · 200 万 · 10%~15%", font_size=Pt(16), font_color=WHITE, alignment=PP_ALIGN.CENTER)
add_line(slide, Inches(0), H-Pt(6), W, Pt(6), PURPLE)


# ══════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════
import os as _os
_os.makedirs("docs/pitch", exist_ok=True)
output_path = "docs/pitch/算力通_IPO退出路径分析.pptx"
prs.save(output_path)
print(f"✅ IPO 退出分析 PPT 已生成: {output_path}")
print(f"   Slides: {len(prs.slides)}")
print("   Size: 16:9 widescreen")
