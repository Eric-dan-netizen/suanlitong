"""
算力通 · 投资人深度分析 PPT 生成器
主题：市场构成 + 开发者趋势 + 退出收益测算
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Design Tokens ────────────────────────────────────
W = Inches(13.333)
H = Inches(7.5)

BG_DARK   = RGBColor(0x0A, 0x0E, 0x17)
BG_CARD   = RGBColor(0x12, 0x18, 0x24)
ACCENT    = RGBColor(0x00, 0xD4, 0xFF)
ACCENT2   = RGBColor(0x00, 0xFF, 0x88)
ACCENT3   = RGBColor(0xFF, 0x6B, 0x35)
ACCENT4   = RGBColor(0xFF, 0xD7, 0x00)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GRAY      = RGBColor(0x88, 0x93, 0xA0)
GRAY_DIM  = RGBColor(0x55, 0x5E, 0x6B)
BORDER    = RGBColor(0x1E, 0x28, 0x36)

FONT_TITLE = "PingFang SC"
FONT_BODY  = "PingFang SC"

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank_layout = prs.slide_layouts[6]

# ── Helpers ──────────────────────────────────────────
def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fill_color=None, border_color=None, border_width=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.fill.solid()
        if border_width:
            shape.line.width = border_width
    return shape

def add_line(slide, l, t, w, h, color, lw=Pt(1)):
    return add_rect(slide, l, t, w, h, fill_color=color)

def add_textbox(slide, l, t, w, h, text="", font_size=Pt(14),
                font_color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                font_name=FONT_BODY, line_spacing=1.2):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.line_spacing = Pt(int(font_size.pt * line_spacing))
    return txBox

def add_ml_text(slide, l, t, w, h, lines, default_size=Pt(13),
                default_color=WHITE, default_bold=False, alignment=PP_ALIGN.LEFT,
                line_spacing=1.3):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(line, str):
            p.text = line
            p.font.size = default_size
            p.font.color.rgb = default_color
            p.font.bold = default_bold
        else:
            txt, sz, clr, bld = line[0], line[1] if len(line)>1 else default_size, \
                                line[2] if len(line)>2 else default_color, \
                                line[3] if len(line)>3 else default_bold
            p.text = txt
            p.font.size = sz
            p.font.color.rgb = clr
            p.font.bold = bld
        p.font.name = FONT_BODY
        p.alignment = alignment
        p.line_spacing = Pt(18)  # fixed line spacing
    return txBox

def add_section_title(slide, title, subtitle=""):
    add_line(slide, Inches(0), Inches(0), W, Pt(4), ACCENT)
    add_textbox(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.55),
                text=title, font_size=Pt(30), font_color=WHITE, bold=True)
    if subtitle:
        add_textbox(slide, Inches(0.6), Inches(0.95), Inches(12), Inches(0.35),
                    text=subtitle, font_size=Pt(14), font_color=GRAY)

def add_card(slide, l, t, w, h, title="", body_lines=None, title_color=ACCENT, body_size=Pt(11)):
    add_rect(slide, l, t, w, h, fill_color=BG_CARD, border_color=BORDER, border_width=Pt(0.5))
    if title:
        add_textbox(slide, l+Inches(0.15), t+Inches(0.1), w-Inches(0.3), Inches(0.28),
                    text=title, font_size=Pt(13), font_color=title_color, bold=True)
    if body_lines:
        add_ml_text(slide, l+Inches(0.15), t+Inches(0.42), w-Inches(0.3), h-Inches(0.5),
                    lines=body_lines, default_size=body_size, default_color=GRAY, line_spacing=1.25)

def add_arrow_right(slide, l, t, color=ACCENT):
    add_textbox(slide, l, t, Inches(0.4), Inches(0.25), text="→", font_size=Pt(18), font_color=color, bold=True)

def add_badge(slide, l, t, w, h, text, bg=ACCENT, text_color=BG_DARK, font_size=Pt(11)):
    add_rect(slide, l, t, w, h, fill_color=bg)
    add_textbox(slide, l, t, w, h, text=text, font_size=font_size, font_color=text_color,
                bold=True, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════
# SLIDE 1 · COVER
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG_DARK)
add_line(slide, Inches(0), Inches(0), W, Pt(6), ACCENT)
add_textbox(slide, Inches(1), Inches(1.8), Inches(11.3), Inches(0.9),
            text="算力通 · 投资人深度分析", font_size=Pt(44), font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_line(slide, Inches(4.5), Inches(2.9), Inches(4.3), Pt(2), ACCENT)
add_ml_text(slide, Inches(1.5), Inches(3.3), Inches(10.3), Inches(1.2),
            lines=[
                ("市场构成拆解  ·  开发者五年趋势  ·  退出收益测算", Pt(18), GRAY, False),
                ("", Pt(10), GRAY, False),
                ("种子轮 · 200 万 · 10%~15%  |  2026 年 7 月", Pt(14), ACCENT2, False),
            ], alignment=PP_ALIGN.CENTER)
add_line(slide, Inches(0), H-Pt(6), W, Pt(6), ACCENT)


# ══════════════════════════════════════════════════════
# SLIDE 2 · 万亿算力市场：四个层次
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG_DARK)
add_section_title(slide, "8,351 亿算力市场，钱花在哪了？",
                  "中国信通院口径：涵盖芯片→服务的全产业链  |  2025 年数据")

layers = [
    ("④ 服务层", "~2,100-2,500亿", "25-30%",
     ["云计算 IaaS/PaaS", "AI 模型训练/推理", "算力调度与运维"],
     ACCENT2, Inches(2.8)),
    ("③ 软件层", "~1,250-1,670亿", "15-20%",
     ["AI 框架 (PyTorch/PP)", "数据库与中间件", "算力管理平台"],
     ACCENT, Inches(1.9)),
    ("② 基础设施层", "~835-1,250亿", "10-15%",
     ["IDC 建设与运营", "电力/冷却/液冷", "1,250万+ 标准机架"],
     ACCENT4, Inches(1.5)),
    ("① 硬件与设备", "~3,750-4,175亿", "45-50%",
     ["GPU/AI 芯片 (含国产)", "服务器整机", "存储与网络设备"],
     ACCENT3, Inches(1.9)),
]

for i, (name, scale, pct, items, color, h) in enumerate(layers):
    y = Inches(1.5 + i * 1.45)
    add_rect(slide, Inches(0.6), y, Inches(5.2), h, fill_color=color)
    add_textbox(slide, Inches(0.7), y+Inches(0.05), Inches(3.5), Inches(0.3),
                text=name, font_size=Pt(16), font_color=BG_DARK, bold=True)
    add_textbox(slide, Inches(3.5), y+Inches(0.05), Inches(2.2), Inches(0.3),
                text=f"{scale}  ({pct})", font_size=Pt(13), font_color=BG_DARK, bold=True, alignment=PP_ALIGN.RIGHT)
    for j, item in enumerate(items):
        add_textbox(slide, Inches(0.9), y+Inches(0.35 + j*0.22), Inches(4.5), Inches(0.2),
                    text=f"• {item}", font_size=Pt(10), font_color=BG_DARK)

# Right side annotations
add_textbox(slide, Inches(6.5), Inches(1.5), Inches(6.2), Inches(0.35),
            text="▸ 算力通的定位", font_size=Pt(16), font_color=ACCENT, bold=True)
add_ml_text(slide, Inches(6.5), Inches(1.95), Inches(6.2), Inches(2.5),
            lines=[
                ("", Pt(6), GRAY, False),
                ("硬件层  3,750亿  →  NVIDIA/华为/浪潮的蛋糕", Pt(12), GRAY, False),
                ("         算力通不碰 ← 不买卡，不造芯片", Pt(11), GRAY_DIM, False),
                ("", Pt(6), GRAY, False),
                ("基础设施  835亿  →  万国数据/世纪互联", Pt(12), GRAY, False),
                ("         二期自有 GPU 才需要 ← 一期不碰", Pt(11), GRAY_DIM, False),
                ("", Pt(6), GRAY, False),
                ("软件层  1,250亿  →  算力通的「差异化武器」", Pt(12), ACCENT, True),
                ("         比价引擎 · 调度系统 · 一键部署", Pt(11), GRAY_DIM, False),
                ("", Pt(6), GRAY, False),
                ("服务层  2,100亿  →  算力通的核心赛道 ★", Pt(12), ACCENT2, True),
                ("         GPU 租赁子市场 ~130-160亿  TAM 97亿", Pt(11), GRAY_DIM, False),
            ])

# Key insight at bottom
add_rect(slide, Inches(0.6), Inches(6.6), Inches(12.1), Inches(0.55), fill_color=BG_CARD, border_color=ACCENT, border_width=Pt(1))
add_textbox(slide, Inches(0.8), Inches(6.65), Inches(11.7), Inches(0.45),
            text="核心洞察：算力通站在服务层与软件层的交叉点 — 不造芯片、不建机房、不和云厂商正面竞争，做云厂商的「体验层」",
            font_size=Pt(13), font_color=WHITE, bold=True)


# ══════════════════════════════════════════════════════
# SLIDE 3 · 按算力类型切分
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG_DARK)
add_section_title(slide, "按算力类型：智能算力一骑绝尘",
                  "2025 年 1,053 EFLOPS，智能算力占比 81%，年增 ~40%")

types = [
    ("智能算力 (AI)", "1,053 EFLOPS", "81%", "+40% / 年", ACCENT2),
    ("通用算力", "~180 EFLOPS", "14%", "平稳增长", ACCENT),
    ("超级算力 (HPC)", "~60 EFLOPS", "5%", "科研/军工", ACCENT4),
]

for i, (name, scale, pct, growth, color) in enumerate(types):
    x = Inches(0.8 + i * 4.1)
    w_bar = Inches(3.6)
    pct_val = float(pct.replace("%",""))
    bar_w = w_bar * (pct_val / 100)

    # Label
    add_textbox(slide, x, Inches(1.6), w_bar, Inches(0.35),
                text=f"{name}", font_size=Pt(16), font_color=color, bold=True)
    add_textbox(slide, x, Inches(2.0), w_bar, Inches(0.25),
                text=f"{scale}  ·  {pct}  ·  {growth}", font_size=Pt(11), font_color=GRAY)
    # Bar
    add_rect(slide, x, Inches(2.4), Inches(bar_w), Inches(0.5), fill_color=BG_CARD)
    add_rect(slide, x, Inches(2.4), Inches(bar_w * pct_val / 100), Inches(0.5), fill_color=color)

# Trend context
add_rect(slide, Inches(0.6), Inches(3.4), Inches(12.1), Inches(3.6), fill_color=BG_CARD, border_color=BORDER)
add_textbox(slide, Inches(0.9), Inches(3.55), Inches(11.5), Inches(0.3),
            text="智能算力增长轨迹 (EFLOPS)", font_size=Pt(15), font_color=ACCENT, bold=True)

years_data = [
    ("2023", "725", "—"),
    ("2024", "—", ""),
    ("2025", "1,053", "+45%"),
    ("2026E", "1,460", "+39%"),
    ("2027E", "~2,000", "+37%"),
    ("2028E", "~2,700", "+35%"),
]

for i, (yr, val, g) in enumerate(years_data):
    x = Inches(1.0 + i * 2.0)
    clr = ACCENT2 if i >= 4 else ACCENT
    add_textbox(slide, x, Inches(4.05), Inches(1.8), Inches(0.28),
                text=yr, font_size=Pt(13), font_color=clr, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x, Inches(4.35), Inches(1.8), Inches(0.28),
                text=val, font_size=Pt(18), font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x, Inches(4.7), Inches(1.8), Inches(0.2),
                text=g, font_size=Pt(10), font_color=GRAY, alignment=PP_ALIGN.CENTER)

# CAGR callout
add_textbox(slide, Inches(1.0), Inches(5.3), Inches(11.3), Inches(0.4),
            text="2023-2028 CAGR: 46.2%   |   日均 Token 调用量两年涨 1,000 倍 (1,000亿 → 140万亿+)",
            font_size=Pt(14), font_color=ACCENT2, bold=True, alignment=PP_ALIGN.CENTER)

# Bottom insight
add_ml_text(slide, Inches(0.9), Inches(5.9), Inches(11.5), Inches(0.9),
            lines=[
                ("智能算力的爆发式增长意味着什么？", Pt(13), WHITE, True),
                ("→ 训练需求之外，推理需求正在成为主力增量 — 每个上线的 AI 应用都在持续消耗 GPU", Pt(12), GRAY, False),
                ("→ 推理需求分散、碎片化、弹性强 — 这恰恰是「按秒计费 + Spot 竞价」模式最擅长的场景", Pt(12), GRAY, False),
            ])


# ══════════════════════════════════════════════════════
# SLIDE 4 · 算力通在产业链的位置
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG_DARK)
add_section_title(slide, "算力通在万亿市场中的位置",
                  "上游芯片→中游云厂商→算力通（聚合层）→下游开发者")

# Value chain visualization
chain = [
    ("芯片厂商", "NVIDIA\n华为昇腾\n寒武纪", ACCENT3, Inches(2.0)),
    ("云厂商", "阿里云 PAI\n华为云\n腾讯云", ACCENT4, Inches(2.0)),
    ("算力通", "多云比价\n一键部署\n按秒计费", ACCENT2, Inches(2.5)),
    ("AI 开发者", "独立开发者\n小团队\nAI 创业者", ACCENT, Inches(2.0)),
]

for i, (name, desc, color, h) in enumerate(chain):
    x = Inches(0.6 + i * 3.2)
    add_rect(slide, x, Inches(1.6), Inches(2.9), h, fill_color=color)
    add_textbox(slide, x+Inches(0.1), Inches(1.65), Inches(2.7), Inches(0.3),
                text=name, font_size=Pt(16), font_color=BG_DARK, bold=True)
    add_ml_text(slide, x+Inches(0.1), Inches(2.0), Inches(2.7), h-Inches(0.5),
                lines=[(desc, Pt(12), BG_DARK, False)])
    if i < 3:
        add_arrow_right(slide, x+Inches(3.0), Inches(2.2), WHITE)

# Layer details below
add_textbox(slide, Inches(0.6), Inches(4.0), Inches(12.1), Inches(0.35),
            text="各层市场规模 & 算力通参与度", font_size=Pt(18), font_color=WHITE, bold=True)

table_data = [
    ("层级", "2025 规模", "算力通参与", "参与方式", "参与阶段"),
    ("硬件", "3,750-4,175亿", "0%", "不参与", "—"),
    ("基础设施", "835-1,250亿", "0% → 5%", "二期租赁 IDC 机柜", "M18+"),
    ("软件", "1,250-1,670亿", "差异化层", "比价引擎/调度 SaaS", "M0+"),
    ("服务·GPU租赁", "130-160亿", "核心赛道 ★", "聚合+返佣+差价", "M3+"),
]

for r, row in enumerate(table_data):
    y = Inches(4.5 + r * 0.42)
    is_header = r == 0
    for c, cell in enumerate(row):
        cols_w = [Inches(1.8), Inches(2.4), Inches(2.2), Inches(3.0), Inches(2.0)]
        x = Inches(0.6) + sum(cw for cw in cols_w[:c])
        if is_header:
            add_rect(slide, x, y, cols_w[c], Inches(0.38), fill_color=ACCENT)
            add_textbox(slide, x+Inches(0.05), y+Inches(0.04), cols_w[c]-Inches(0.1), Inches(0.3),
                        text=cell, font_size=Pt(11), font_color=BG_DARK, bold=True, alignment=PP_ALIGN.CENTER)
        else:
            bg = BG_CARD if r % 2 == 0 else BG_DARK
            txt_clr = ACCENT2 if cell and "★" in cell else WHITE
            add_rect(slide, x, y, cols_w[c], Inches(0.38), fill_color=bg, border_color=BORDER, border_width=Pt(0.3))
            add_textbox(slide, x+Inches(0.05), y+Inches(0.04), cols_w[c]-Inches(0.1), Inches(0.3),
                        text=cell, font_size=Pt(10), font_color=txt_clr, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════
# SLIDE 5 · AI 开发者五年趋势总览
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG_DARK)
add_section_title(slide, "AI 开发者五年趋势：总量翻倍 + 结构重组",
                  "2025 年 940 万人 → 2030 年 ~2,200 万人  |  CAGR ~18-20%")

# Year-by-year growth bar
years = [("2025", 940, "基准"), ("2026E", 1160, "+23%"), ("2027E", 1400, "+21%"),
         ("2028E", 1670, "+19%"), ("2029E", 1970, "+18%"), ("2030E", 2200, "+12%")]

for i, (yr, val, g) in enumerate(years):
    x = Inches(0.6 + i * 2.1)
    bar_h = Inches(val / 940 * 2.5)  # proportional
    add_rect(slide, x+Inches(0.35), Inches(3.3 - bar_h.emu/Inches(1).emu), Inches(1.4), bar_h,
             fill_color=ACCENT2 if i >= 3 else ACCENT)
    add_textbox(slide, x, Inches(3.4), Inches(2.1), Inches(0.28),
                text=f"{val}万", font_size=Pt(18) if val>=1400 else Pt(16), font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x, Inches(3.7), Inches(2.1), Inches(0.2),
                text=f"{yr}  {g}", font_size=Pt(10), font_color=GRAY if i>0 else ACCENT2, alignment=PP_ALIGN.CENTER)

# Driving forces
add_textbox(slide, Inches(0.6), Inches(4.2), Inches(12.1), Inches(0.35),
            text="五大增长引擎", font_size=Pt(16), font_color=ACCENT, bold=True)

drivers = [
    ("AI 编程工具", "Cursor/Copilot 让\n1人=1个团队", ACCENT2),
    ("开源模型", "DeepSeek/Qwen\n逼近 GPT-4 水平", ACCENT),
    ("AI Agent 创业潮", "每个 Agent = 1个\n微型 SaaS 产品", ACCENT4),
    ("Token 价格暴跌", "价格降 90%\n用量涨 10 倍+", ACCENT3),
    ("AI 教育普及", "500+ 高校开设\nAI 专业课程", WHITE),
]

for i, (title, desc, color) in enumerate(drivers):
    x = Inches(0.6 + i * 2.5)
    add_card(slide, x, Inches(4.65), Inches(2.2), Inches(1.5),
             title=title, body_lines=[desc], title_color=color, body_size=Pt(11))

# Bottom insight
add_rect(slide, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.65), fill_color=BG_CARD, border_color=ACCENT2)
add_ml_text(slide, Inches(0.9), Inches(6.55), Inches(11.5), Inches(0.55),
            lines=[
                ("先行指标印证：Token 调用量两年 1,000 倍 → AI 核心产业 1.2 万亿 →「十五五」末 10 万亿", Pt(13), WHITE, True),
                ("开发者增速 (18-22%) < 算力增速 (46%) — 因为人比机器慢，但 AI 工具正在加速这个进程", Pt(11), GRAY, False),
            ])


# ══════════════════════════════════════════════════════
# SLIDE 6 · 结构性迁移：四类人群此消彼长
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG_DARK)
add_section_title(slide, "结构性迁移：饼在变大，更在重新切",
                  "独立开发者 + 小团队占比从 34% → 47% —「超级个体化」是最大的结构红利")

segments = [
    ("独立开发者", "129万→480万", "+272%", "13.7%→22%", ACCENT2, "涨幅最大\nAI工具让个人=\n一个团队"),
    ("50人以下小团队", "195万→550万", "+182%", "20.7%→25%", ACCENT, "创业潮主力\n3-5人即可\n做出 AI 产品"),
    ("千人以上大厂", "228万→320万", "+40%", "24.3%→14.5%", ACCENT4, "绝对量涨\n相对量降\nAI替代传统岗位"),
    ("其他(学生/科研等)", "388万→850万", "+119%", "41.3%→38.5%", WHITE, "高校扩招\n科研向AI迁移\n中型企业转型"),
]

for i, (name, scale, growth, pct_change, color, insight) in enumerate(segments):
    x = Inches(0.5 + i * 3.15)
    add_rect(slide, x, Inches(1.5), Inches(2.9), Inches(0.45), fill_color=color)
    add_textbox(slide, x+Inches(0.1), Inches(1.52), Inches(2.7), Inches(0.4),
                text=name, font_size=Pt(14), font_color=BG_DARK, bold=True)
    add_ml_text(slide, x, Inches(2.1), Inches(2.9), Inches(2.0),
                lines=[
                    (scale, Pt(18), color, True),
                    (growth, Pt(13), ACCENT2 if "+272" in growth else GRAY, True),
                    (f"占比 {pct_change}", Pt(11), GRAY, False),
                    ("", Pt(6), GRAY, False),
                    (insight, Pt(10), GRAY, False),
                ])

# The "super individual" callout
add_rect(slide, Inches(0.6), Inches(4.5), Inches(12.1), Inches(2.5), fill_color=BG_CARD, border_color=ACCENT2, border_width=Pt(1))
add_textbox(slide, Inches(0.9), Inches(4.6), Inches(11.5), Inches(0.35),
            text="🔥 「超级个体化」—— 算力通最大的结构性红利", font_size=Pt(18), font_color=ACCENT2, bold=True)

add_ml_text(slide, Inches(0.9), Inches(5.1), Inches(5.5), Inches(1.7),
            lines=[
                ("五个叠加趋势：", Pt(13), WHITE, True),
                ("1. Cursor/Claude Code — 1 个全栈 = 去年 5 人团队", Pt(12), GRAY, False),
                ("2. Qwen/DeepSeek 开源 — 微调成本从百万降至数千", Pt(12), GRAY, False),
                ("3. AI Agent — 不需要「做大平台」，做「一个能卖钱的 Agent」", Pt(12), GRAY, False),
                ("4. Token 价格逼近一度电 — 推理成本趋零", Pt(12), GRAY, False),
                ("5. 小程序/抖音 AI 能力开放 — 独立开发者有现成分发渠道", Pt(12), GRAY, False),
            ])

add_ml_text(slide, Inches(6.8), Inches(5.1), Inches(5.5), Inches(1.7),
            lines=[
                ("对算力通意味着什么？", Pt(13), ACCENT2, True),
                ("", Pt(6), GRAY, False),
                ("→ 129 万独立开发者不是终点，是起点", Pt(12), GRAY, False),
                ("→ 每年新增 40-60 万独立 AI 开发者", Pt(12), GRAY, False),
                ("→ 每个人都需要 GPU（微调/推理/部署）", Pt(12), GRAY, False),
                ("→ 每个人都不愿被包月锁死", Pt(12), GRAY, False),
                ("→ 每个人都不想花 2 天配环境", Pt(12), GRAY, False),
                ("", Pt(8), GRAY, False),
                ("这就是算力通「按秒计费 + 一键部署」", Pt(13), ACCENT2, True),
                ("天然适配超级个体时代的理由。", Pt(13), ACCENT2, True),
            ])


# ══════════════════════════════════════════════════════
# SLIDE 7 · SAM 增长测算
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG_DARK)
add_section_title(slide, "算力通 SAM：324 万 → 1,030 万（3.2 倍）",
                  "总量翻 2.3 倍 × 核心人群占比升至 47% × 需 GPU 比例升至 70%")

# SAM growth table
sam_metrics = [
    ("指标", "2025", "2030E", "倍数"),
    ("总 AI 开发者", "940 万", "2,200 万", "2.3x"),
    ("独立+小团队 (核心 SAM)", "324 万", "1,030 万", "3.2x"),
    ("其中需 GPU 算力", "194 万 (60%)", "720 万 (70%)", "3.7x"),
]

for r, row in enumerate(sam_metrics):
    y = Inches(1.6 + r * 0.55)
    is_header = r == 0
    cols_w = [Inches(3.5), Inches(2.8), Inches(2.8), Inches(2.8)]
    for c, cell in enumerate(row):
        x = Inches(0.6) + sum(cw for cw in cols_w[:c])
        if is_header:
            add_rect(slide, x, y, cols_w[c], Inches(0.5), fill_color=ACCENT)
            add_textbox(slide, x, y+Inches(0.08), cols_w[c], Inches(0.35),
                        text=cell, font_size=Pt(14), font_color=BG_DARK, bold=True, alignment=PP_ALIGN.CENTER)
        else:
            bg = BG_CARD if r == 2 else BG_DARK
            txt_clr = ACCENT2 if cell and "3." in cell else WHITE
            add_rect(slide, x, y, cols_w[c], Inches(0.5), fill_color=bg, border_color=BORDER, border_width=Pt(0.3))
            add_textbox(slide, x, y+Inches(0.08), cols_w[c], Inches(0.35),
                        text=cell, font_size=Pt(16) if r==2 else Pt(14), font_color=txt_clr, bold=True, alignment=PP_ALIGN.CENTER)

# "为什么需要 GPU 的比例会上升" box
add_rect(slide, Inches(0.6), Inches(3.9), Inches(5.8), Inches(2.5), fill_color=BG_CARD, border_color=ACCENT)
add_textbox(slide, Inches(0.9), Inches(4.0), Inches(5.2), Inches(0.35),
            text="为什么「需 GPU」比例从 60% → 70%？", font_size=Pt(15), font_color=ACCENT, bold=True)
add_ml_text(slide, Inches(0.9), Inches(4.45), Inches(5.2), Inches(1.8),
            lines=[
                ("• AI 正从「大模型公司专属」变成", Pt(12), GRAY, False),
                ("  「所有软件的基础设施」", Pt(12), GRAY, False),
                ("• 五年后，不做 AI 功能的 App", Pt(12), GRAY, False),
                ("  就像今天不联网的 App 一样罕见", Pt(12), GRAY, False),
                ("• 前端/UI/后端开发者都将被拉入", Pt(12), GRAY, False),
                ("  AI 开发 — Cursor 让非 ML 工程师", Pt(12), GRAY, False),
                ("  也能调模型、做微调、跑推理", Pt(12), GRAY, False),
            ])

# Per-capita GPU usage
add_rect(slide, Inches(6.8), Inches(3.9), Inches(5.9), Inches(2.5), fill_color=BG_CARD, border_color=ACCENT2)
add_textbox(slide, Inches(7.1), Inches(4.0), Inches(5.3), Inches(0.35),
            text="人均 GPU 消费：涨还是跌？", font_size=Pt(15), font_color=ACCENT2, bold=True)

add_ml_text(slide, Inches(7.1), Inches(4.45), Inches(5.3), Inches(1.8),
            lines=[
                ("↓  推理效率提升 → 同样效果只需 1/10 算力", Pt(11), GRAY, False),
                ("↑  AI Agent 调用链变长 → 一次任务调 10-50 次", Pt(11), ACCENT2, True),
                ("↑  Token 价格暴跌 → 用量弹性 > 1（Jevons 悖论）", Pt(11), ACCENT2, True),
                ("", Pt(6), GRAY, False),
                ("净效应：人均 GPU 消费持续上升", Pt(13), ACCENT2, True),
                ("证据：Token 价格降了 90%+，但日均调用量", Pt(11), GRAY, False),
                ("从 1,000 亿涨到 140 万亿 — 涨了 1,000 倍", Pt(11), GRAY, False),
            ])

# Bottom summary
add_rect(slide, Inches(0.6), Inches(6.7), Inches(12.1), Inches(0.5), fill_color=BG_CARD, border_color=ACCENT2)
add_textbox(slide, Inches(0.9), Inches(6.75), Inches(11.5), Inches(0.4),
            text="SAM 3.2 倍 × 人均 GPU 消费 1.5-2 倍 → 可服务市场从 50 亿 → 200-300 亿  |  这正是 A 轮窗口（M18-M24）的逻辑基础",
            font_size=Pt(13), font_color=WHITE, bold=True)


# ══════════════════════════════════════════════════════
# SLIDE 8 · 退出收益测算 — 假设前提
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG_DARK)
add_section_title(slide, "退出收益测算 · 建模假设",
                  "所有参数可溯源至商业计划书 V3  |  偏差项用行业惯例补齐")

assumptions = [
    ("种子轮投资额", "200 万", "BP 明确"),
    ("种子轮股权", "12.5%（10%-15% 中值）", "BP 明确"),
    ("期权池", "15%（种子轮）→ 20%（A 轮前）", "BP + 行业惯例"),
    ("A 轮融资", "2,000 万（1,500-2,500 万中值）", "BP 明确 · M18-M20"),
    ("A 轮投前估值", "~8,600 万（ARR ~720万 × 12x）", "BP 参考 + SaaS 中值"),
    ("B 轮", "3 年退出场景下不触发", "BP: B轮在 2028-2029"),
    ("退出倍数", "ARR 6x-9x（中值 7.5x）", "收窄自 BP 5-10x"),
    ("三年 ARR 三情景", "保守 1,620 / 基准 3,600 / 乐观 8,640 万", "BP 5.4 节"),
]

for i, (param, value, source) in enumerate(assumptions):
    y = Inches(1.5 + i * 0.62)
    add_rect(slide, Inches(0.6), y, Inches(12.1), Inches(0.55), fill_color=BG_CARD, border_color=BORDER, border_width=Pt(0.3))
    add_textbox(slide, Inches(0.8), y+Inches(0.08), Inches(3.5), Inches(0.4),
                text=param, font_size=Pt(13), font_color=WHITE, bold=True)
    add_textbox(slide, Inches(4.5), y+Inches(0.08), Inches(4.5), Inches(0.4),
                text=value, font_size=Pt(13), font_color=ACCENT2, bold=True)
    add_textbox(slide, Inches(9.2), y+Inches(0.08), Inches(3.3), Inches(0.4),
                text=source, font_size=Pt(10), font_color=GRAY_DIM, alignment=PP_ALIGN.RIGHT)

add_rect(slide, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.6), fill_color=BG_CARD, border_color=ACCENT)
add_textbox(slide, Inches(0.9), Inches(6.55), Inches(11.5), Inches(0.5),
            text="核心差异：BP 中的「9-27 倍」是稀释前理论值。本测算计入 A 轮稀释（期权池扩 + A 轮新发），推算投资人实际到手回报。",
            font_size=Pt(13), font_color=WHITE, bold=True)


# ══════════════════════════════════════════════════════
# SLIDE 9 · 股权稀释路径
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG_DARK)
add_section_title(slide, "股权稀释路径：200 万买的 12.5%，三年后还剩多少？",
                  "种子 → 期权池扩张 → A 轮 → 退出  |  累计稀释 ~24%")

# Cap table evolution
dilution_steps = [
    ("M0 · 种子轮 close", "12.50%", "—", "200 万买入，期权池 15% 从创始人出", ACCENT2),
    ("M18 · 期权池 15%→20%", "11.76%", "-5.9%", "扩池 5%，种/创按比例稀释", ACCENT),
    ("M18 · A 轮 close", "9.55%", "-18.8%", "A 轮新发 ~18.8%，全体稀释", ACCENT4),
    ("M36 · 退出", "9.55%", "—", "三年期满，最终持股", ACCENT2),
]

for i, (event, stake, dilution, note, color) in enumerate(dilution_steps):
    y = Inches(1.5 + i * 0.9)
    # Timeline node
    add_rect(slide, Inches(0.6), y+Inches(0.15), Inches(0.25), Inches(0.25), fill_color=color)
    if i < 3:
        add_line(slide, Inches(0.7), y+Inches(0.38), Pt(2), Inches(0.55), color)
    # Event
    add_textbox(slide, Inches(1.1), y, Inches(3.5), Inches(0.3),
                text=event, font_size=Pt(14), font_color=WHITE, bold=True)
    # Stake badge
    add_rect(slide, Inches(4.8), y, Inches(1.5), Inches(0.5), fill_color=color)
    add_textbox(slide, Inches(4.8), y+Inches(0.06), Inches(1.5), Inches(0.38),
                text=stake, font_size=Pt(20), font_color=BG_DARK, bold=True, alignment=PP_ALIGN.CENTER)
    # Dilution
    add_textbox(slide, Inches(6.5), y+Inches(0.06), Inches(1.2), Inches(0.38),
                text=dilution, font_size=Pt(14), font_color=ACCENT3, bold=True)
    # Note
    add_textbox(slide, Inches(7.8), y+Inches(0.06), Inches(5.0), Inches(0.38),
                text=note, font_size=Pt(11), font_color=GRAY)

# Summary box
add_rect(slide, Inches(0.6), Inches(5.2), Inches(5.8), Inches(1.8), fill_color=BG_CARD, border_color=ACCENT2)
add_textbox(slide, Inches(0.9), Inches(5.3), Inches(5.2), Inches(0.3),
            text="稀释总结", font_size=Pt(16), font_color=ACCENT2, bold=True)
add_ml_text(slide, Inches(0.9), Inches(5.7), Inches(5.2), Inches(1.2),
            lines=[
                ("种子轮 12.5% → 退出时 ~9.55%", Pt(16), WHITE, True),
                ("名义股权打 76 折", Pt(14), GRAY, False),
                ("行业正常水平：种子→退出 20-30% 稀释", Pt(12), GRAY_DIM, False),
                ("属于标准区间", Pt(12), GRAY_DIM, False),
            ])

# B-round scenario
add_rect(slide, Inches(6.8), Inches(5.2), Inches(5.9), Inches(1.8), fill_color=BG_CARD, border_color=ACCENT3)
add_textbox(slide, Inches(7.1), Inches(5.3), Inches(5.3), Inches(0.3),
            text="若触发 B 轮（5 年退出场景）", font_size=Pt(16), font_color=ACCENT3, bold=True)
add_ml_text(slide, Inches(7.1), Inches(5.7), Inches(5.3), Inches(1.2),
            lines=[
                ("B 轮 5,000 万-1 亿  →  额外稀释 ~15-20%", Pt(14), WHITE, True),
                ("最终持股 ~6.5-7.5%", Pt(14), GRAY, False),
                ("但退出估值可能翻倍（ARR 过亿）", Pt(12), ACCENT2, True),
                ("净回报通常更高 — 饼更大", Pt(12), ACCENT2, True),
            ])


# ══════════════════════════════════════════════════════
# SLIDE 10 · 全情景回报矩阵
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG_DARK)
add_section_title(slide, "全情景退出回报矩阵",
                  "200 万投入 → 种子投资人实际到手（含 A 轮稀释至 9.55%）")

# Main return table
add_textbox(slide, Inches(0.6), Inches(1.4), Inches(12.1), Inches(0.35),
            text="退出估值 = M36 ARR × 倍数  |  种子回报 = 退出估值 × 9.55%",
            font_size=Pt(13), font_color=GRAY, bold=True)

scenarios = [
    ("情景", "ARR (万)", "退出倍数", "退出估值", "种子回报", "MOIC", "3年 IRR"),
    ("保守·低", "1,620", "5x", "0.81 亿", "774 万", "3.9x", "57%"),
    ("保守·高", "1,620", "7x", "1.13 亿", "1,079 万", "5.4x", "75%"),
    ("基准·低", "3,600", "6x", "2.16 亿", "2,063 万", "10.3x", "117%"),
    ("基准·中 ★", "3,600", "7.5x", "2.70 亿", "2,579 万", "12.9x", "134%"),
    ("基准·高", "3,600", "9x", "3.24 亿", "3,094 万", "15.4x", "149%"),
    ("乐观·低", "8,640", "8x", "6.91 亿", "6,599 万", "33.0x", "220%"),
    ("乐观·高", "8,640", "10x", "8.64 亿", "8,251 万", "41.3x", "245%"),
]

for r, row in enumerate(scenarios):
    y = Inches(1.85 + r * 0.52)
    is_header = r == 0
    cols_w = [Inches(1.5), Inches(1.5), Inches(1.5), Inches(1.8), Inches(2.0), Inches(1.4), Inches(1.6)]
    for c, cell in enumerate(row):
        x = Inches(0.6) + sum(cw for cw in cols_w[:c])
        if is_header:
            add_rect(slide, x, y, cols_w[c], Inches(0.45), fill_color=ACCENT)
            add_textbox(slide, x, y+Inches(0.06), cols_w[c], Inches(0.33),
                        text=cell, font_size=Pt(11), font_color=BG_DARK, bold=True, alignment=PP_ALIGN.CENTER)
        else:
            is_highlight = "★" in cell
            bg = BG_DARK if r % 2 == 0 else BG_CARD
            if is_highlight:
                bg = RGBColor(0x0A, 0x28, 0x20)  # dark green tint
            txt_clr = ACCENT2 if (is_highlight or ("x" in cell and "万" not in cell)) else WHITE
            bld = is_highlight or (c >= 5)
            add_rect(slide, x, y, cols_w[c], Inches(0.45), fill_color=bg, border_color=BORDER, border_width=Pt(0.3))
            add_textbox(slide, x, y+Inches(0.06), cols_w[c], Inches(0.33),
                        text=cell, font_size=Pt(12) if is_highlight else Pt(11), font_color=txt_clr, bold=bld, alignment=PP_ALIGN.CENTER)

# Key takeaway
add_rect(slide, Inches(0.6), Inches(5.75), Inches(12.1), Inches(1.4), fill_color=BG_CARD, border_color=ACCENT2)
add_ml_text(slide, Inches(0.9), Inches(5.85), Inches(11.5), Inches(1.2),
            lines=[
                ("基准方案（最可能）：ARR 3,600 万 → 退出估值 2.70 亿 → 种子回报 2,579 万 → 12.9x MOIC", Pt(16), ACCENT2, True),
                ("", Pt(6), GRAY, False),
                ("保守方案底线：ARR 1,620 万 + 最低倍数 → 仍有 3.9x MOIC → 拿回近 4 倍，远高于「亏钱」线", Pt(13), WHITE, False),
                ("乐观方案上限：ARR 8,640 万 + 战略溢价 → 41x MOIC → 200 万变 8,200 万，基金回报级项目", Pt(13), ACCENT2, False),
            ])


# ══════════════════════════════════════════════════════
# SLIDE 11 · 三个关键阈值
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG_DARK)
add_section_title(slide, "投资人视角：三个关键阈值",
                  "回本线  ·  目标线  ·  超额线")

thresholds = [
    ("🟢 回本线", "ARR ~280-420 万", "不到 M12 目标即可回本",
     "退出估值 = 200 万 ÷ 9.55% = 2,094 万\nARR = 2,094 ÷ 5~7.5x = 279~419 万\n\n→ M6-M10 之间就能跨过\n→ 下行风险极低，本金安全垫厚", ACCENT2, Inches(2.5)),
    ("🔵 目标线 · 10x MOIC", "ARR ~2,800 万", "VC 种子轮及格门槛",
     "10x MOIC = 2,000 万 ÷ 9.55% = 2.09 亿\nARR = 2.09 亿 ÷ 7.5x = 2,793 万\n\n→ 基准方案 3,600 万完全覆盖\n→ 按计划执行 = 投资人满意回报", ACCENT, Inches(2.5)),
    ("🟣 超额线 · 20x MOIC", "ARR ~5,600 万", "基金回报级项目",
     "20x MOIC = 4,000 万 ÷ 9.55% = 4.19 亿\nARR = 4.19 亿 ÷ 7.5x = 5,587 万\n\n→ 乐观方案 8,640 万轻松超过\n→ 基准方案接近但差一点", ACCENT3, Inches(2.5)),
]

for i, (title, arr_req, tagline, detail, color, h) in enumerate(thresholds):
    x = Inches(0.5 + i * 4.2)
    add_rect(slide, x, Inches(1.5), Inches(3.9), Inches(0.5), fill_color=color)
    add_textbox(slide, x+Inches(0.1), Inches(1.53), Inches(3.7), Inches(0.44),
                text=title, font_size=Pt(18), font_color=BG_DARK, bold=True)
    add_textbox(slide, x, Inches(2.1), Inches(3.9), Inches(0.25),
                text=arr_req, font_size=Pt(14), font_color=ACCENT2, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x, Inches(2.4), Inches(3.9), Inches(0.2),
                text=tagline, font_size=Pt(11), font_color=GRAY, alignment=PP_ALIGN.CENTER)
    add_card(slide, x, Inches(2.7), Inches(3.9), h,
             title="", body_lines=[detail], body_size=Pt(11))

# Benchmark comparison
add_textbox(slide, Inches(0.6), Inches(5.5), Inches(12.1), Inches(0.35),
            text="算力通 vs 典型种子轮 SaaS 项目", font_size=Pt(18), font_color=WHITE, bold=True)

bench_cols = ["指标", "典型 SaaS 种子轮", "算力通(基准,含稀释)", "优势来源"]
bench_data = [
    ("3 年 MOIC 中位数", "3-5x", "12.9x", "毛利率结构好"),
    ("3 年 IRR 中位数", "40-60%", "134%", "资本效率高 LTV/CAC=16:1"),
    ("回本概率", "~60%", ">95%", "ARR 仅需 280 万即回本"),
    ("10x 概率", "~15-20%", ">50%", "ARR 2,800 万即可"),
]

for r, row in enumerate([bench_cols] + bench_data):
    y = Inches(5.95 + r * 0.4)
    is_header = r == 0
    bw = [Inches(2.8), Inches(2.9), Inches(3.3), Inches(3.1)]
    for c, cell in enumerate(row):
        x = Inches(0.6) + sum(cw for cw in bw[:c])
        if is_header:
            add_rect(slide, x, y, bw[c], Inches(0.36), fill_color=ACCENT)
            add_textbox(slide, x, y+Inches(0.03), bw[c], Inches(0.3),
                        text=cell, font_size=Pt(10), font_color=BG_DARK, bold=True, alignment=PP_ALIGN.CENTER)
        else:
            bg = BG_DARK if r % 2 == 0 else BG_CARD
            txt_clr = ACCENT2 if r == 1 and c == 2 else WHITE
            add_rect(slide, x, y, bw[c], Inches(0.36), fill_color=bg, border_color=BORDER, border_width=Pt(0.3))
            add_textbox(slide, x, y+Inches(0.03), bw[c], Inches(0.3),
                        text=cell, font_size=Pt(10), font_color=txt_clr, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════
# SLIDE 12 · 战略收购的额外上行
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG_DARK)
add_section_title(slide, "最大变量不是 ARR，是「谁买」",
                  "财务退出 (ARR × 倍数) 是地面  ·  战略收购是天花板")

add_ml_text(slide, Inches(0.6), Inches(1.5), Inches(12.1), Inches(1.2),
            lines=[
                ("两种估值逻辑的差异", Pt(18), WHITE, True),
                ("", Pt(6), GRAY, False),
                ("财务退出：估值 = ARR × 倍数 → 基于当下收入", Pt(14), GRAY, False),
                ("战略收购：估值 = 用户价值 × 战略溢价 → 基于「买下你能省我多少时间/钱」", Pt(14), ACCENT2, False),
            ])

# Strategic acquirers
acquirers = [
    ("阿里云", "缺「个人开发者」用户层\n算力通的 3,000 付费用户\n是 PAI 触达不到的长尾", "→ 年收入是算力通 SAM 200 倍\n→ 为一个功能改造后台不划算\n→ 买下来最省事", ACCENT),
    ("字节·火山引擎", "正在追赶阿里云/华为云\n需要现成的用户入口\n比价引擎 =「去哪儿」模式", "→ 2015 携程合并去哪儿\n→ 比价平台被上游收购\n→ 是最经典的互联网退出路径", ACCENT4),
    ("华为云", "昇腾芯片正在铺量\n需要线上销售渠道\n算力通 = 昇腾的「天猫」", "→ 国产替代是国家级战略\n→ 渠道比芯片更稀缺\n→ 战略价值可能远超财务估值", ACCENT3),
]

for i, (name, logic, insight, color) in enumerate(acquirers):
    x = Inches(0.4 + i * 4.25)
    add_rect(slide, x, Inches(2.9), Inches(4.0), Inches(0.45), fill_color=color)
    add_textbox(slide, x+Inches(0.1), Inches(2.93), Inches(3.8), Inches(0.4),
                text=name, font_size=Pt(16), font_color=BG_DARK, bold=True)
    add_card(slide, x, Inches(3.45), Inches(4.0), Inches(1.4),
             title="", body_lines=[logic], body_size=Pt(11))
    add_card(slide, x, Inches(4.95), Inches(4.0), Inches(1.2),
             title="", body_lines=[insight], body_size=Pt(10))

# Scenario comparison
add_rect(slide, Inches(0.6), Inches(6.3), Inches(12.1), Inches(0.75), fill_color=BG_CARD, border_color=ACCENT2, border_width=Pt(1))
add_ml_text(slide, Inches(0.9), Inches(6.35), Inches(11.5), Inches(0.65),
            lines=[
                ("财务退出 12.9x MOIC 是保底  →  战略收购可轻松翻到 20x+", Pt(18), ACCENT2, True),
                ("但战略收购可遇不可求 — 不能作为融资假设，只能作为上行惊喜。", Pt(12), GRAY, False),
            ])


# ══════════════════════════════════════════════════════
# SLIDE 13 · 风险调整后结论
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG_DARK)
add_section_title(slide, "风险调整后的结论",
                  "算力通的回报特征在种子轮项目中属于优质档位")

conclusions = [
    ("下行保护极强", "ARR 仅需 280 万即回本\n保守方案仍 3.9x MOIC\n本金亏损概率 < 5%", ACCENT2),
    ("基准回报达标", "12.9x MOIC 超过 VC 10x 目标\nIRR 134% 远超替代投资\nLTV/CAC = 16:1 模型健康", ACCENT),
    ("上行空间可观", "乐观方案 41x MOIC\n战略收购可能翻倍\nSAM 3.2 倍增长提供顺风", ACCENT3),
]

for i, (title, detail, color) in enumerate(conclusions):
    x = Inches(0.5 + i * 4.25)
    add_rect(slide, x, Inches(1.5), Inches(4.0), Inches(0.5), fill_color=color)
    add_textbox(slide, x+Inches(0.1), Inches(1.53), Inches(3.8), Inches(0.44),
                text=title, font_size=Pt(18), font_color=BG_DARK, bold=True)
    add_card(slide, x, Inches(2.1), Inches(4.0), Inches(1.8),
             title="", body_lines=[detail], body_size=Pt(12))

# Three reasons
add_rect(slide, Inches(0.6), Inches(4.2), Inches(12.1), Inches(2.8), fill_color=BG_CARD, border_color=ACCENT)
add_textbox(slide, Inches(0.9), Inches(4.3), Inches(11.5), Inches(0.35),
            text="为什么算力通比大多数种子轮 SaaS 更安全？三个「低」", font_size=Pt(18), font_color=ACCENT, bold=True)

reasons = [
    ("回本线极低", "代理模式不烧钱\n每个用户都是正向现金流\nM6-M10 即可盈亏平衡", "ARR 280 万即回本 vs 典型 500-1,000 万"),
    ("资本效率高", "LTV/CAC = 16:1\n不需要烧钱换增长\n内容营销 + 社区驱动获客", "行业健康线 LTV/CAC ≥ 3:1\n算力通是 5 倍于健康线"),
    ("退出倍数有锚", "CoreWeave 12x ARR IPO\nPaperspace 8x 被收购\n算力平台估值有公开对标", "纯 SaaS 估值靠「讲增长故事」\n算力通有硬资产 + 可比交易"),
]

for i, (title, detail, compare) in enumerate(reasons):
    y = Inches(4.75 + i * 0.7)
    add_textbox(slide, Inches(1.0), y, Inches(2.0), Inches(0.3),
                text=title, font_size=Pt(15), font_color=ACCENT2, bold=True)
    add_textbox(slide, Inches(3.2), y, Inches(4.5), Inches(0.3),
                text=detail, font_size=Pt(12), font_color=WHITE, bold=False)
    add_textbox(slide, Inches(7.8), y, Inches(4.5), Inches(0.3),
                text=compare, font_size=Pt(11), font_color=GRAY, bold=False)


# ══════════════════════════════════════════════════════
# SLIDE 14 · END PAGE
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG_DARK)
add_line(slide, Inches(0), Inches(0), W, Pt(6), ACCENT)

add_textbox(slide, Inches(1.5), Inches(1.5), Inches(10.3), Inches(0.9),
            text="算力通 · 投资人深度分析", font_size=Pt(42), font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_line(slide, Inches(4.5), Inches(2.6), Inches(4.3), Pt(2), ACCENT)

add_ml_text(slide, Inches(1.5), Inches(3.0), Inches(10.3), Inches(1.8),
            lines=[
                ("市场构成  ·  开发者趋势  ·  退出收益", Pt(20), GRAY, False),
                ("", Pt(10), GRAY, False),
                ("本报告基于商业计划书 V3 数据与公开行业资料", Pt(14), GRAY_DIM, False),
                ("不构成投资建议  ·  仅供种子轮投资人审阅", Pt(14), GRAY_DIM, False),
                ("", Pt(10), GRAY, False),
                ("算力通 · 让算力像自来水一样，扭开即用", Pt(18), ACCENT2, True),
            ], alignment=PP_ALIGN.CENTER)

add_line(slide, Inches(5.5), Inches(5.2), Inches(2.3), Pt(2), ACCENT)

add_textbox(slide, Inches(1.5), Inches(5.5), Inches(10.3), Inches(0.5),
            text="创始人：Eric  |  种子轮 · 200 万 · 10%~15%", font_size=Pt(16), font_color=WHITE, alignment=PP_ALIGN.CENTER)

add_line(slide, Inches(0), H-Pt(6), W, Pt(6), ACCENT)


# ══════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════
import os as _os
_os.makedirs("docs/pitch", exist_ok=True)
output_path = "docs/pitch/算力通_投资人深度分析.pptx"
prs.save(output_path)
print(f"✅ 投资人深度分析 PPT 已生成: {output_path}")
print(f"   Slides: {len(prs.slides)}")
print("   Size: 16:9 widescreen (13.33\" × 7.5\")")
